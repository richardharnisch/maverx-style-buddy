#!/usr/bin/env python3
"""Build one .pptx by cloning master template slides and swapping text.

Usage:
    python build_deck.py --session-plan sessions/1/session_plan.json \
                         --master assets/maverx_master.pptx \
                         --catalog assets/template_catalog.json \
                         --out sessions/1/session_1.pptx
"""
from __future__ import annotations

import argparse
import copy
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from lxml import etree

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent

# Map slide.role -> catalog key
ROLE_TO_CATALOG = {
    "cover": "cover",
    "agenda": "agenda",
    "about_session": "about_session",
    "timetable": "timetable",
    "section_divider_kickoff": ("section_divider", "kickoff"),
    "section_divider_theory": ("section_divider", "theory"),
    "section_divider_exercise": ("section_divider", "exercise"),
    "section_divider_wrapup": ("section_divider", "wrapup"),
    "theory_text": "theory_text",
    "theory_two_column": "theory_two_column",
    "theory_bullets": "theory_bullets",
    "theory_visual": "theory_visual",
    "example": "example",
    "exercise_brief": "exercise_brief",
    "exercise_steps": "exercise_steps",
    "debrief": "debrief",
    "wrapup_takeaways": "wrapup_takeaways",
    "wrapup_next": "wrapup_next",
    "resources": "resources",
    "closer": "closer",
    "break": "break",
    "mentimeter_recap": "mentimeter_recap",
    "asset_template": "asset_template",
    "default_text": "default_text",
}


def resolve_template_idx(catalog: dict, role: str) -> int:
    key = ROLE_TO_CATALOG.get(role, "default_text")
    if isinstance(key, tuple):
        return int(catalog[key[0]][key[1]])
    val = catalog.get(key, catalog["default_text"])
    if isinstance(val, dict):
        # Should not happen with our mapping
        val = next(iter(val.values()))
    return int(val)


def clone_slide(prs: Presentation, src_slide) -> object:
    """Duplicate src_slide into prs, returning the new slide. Preserves all shapes/notes layout."""
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # Remove any default placeholders the new slide inherited
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    # Copy every shape from src
    for shp in src_slide.shapes:
        el = shp.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide


def iter_text_frames(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame


def replace_text_in_frame(tf, new_text: str) -> None:
    """Replace all text in a TextFrame with new_text, keeping the formatting of the first run."""
    # Wipe all paragraphs except the first; keep first paragraph's first run formatting.
    paras = list(tf.paragraphs)
    if not paras:
        tf.text = new_text
        return
    # collapse to one paragraph
    first = paras[0]
    # remove extra paragraphs
    for p in paras[1:]:
        p._p.getparent().remove(p._p)
    # remove runs after first, keep first run's formatting
    runs = list(first.runs)
    if not runs:
        first.text = new_text
        return
    first.text = new_text  # keeps run0 formatting in python-pptx
    # remove leftover runs created by first.text reset
    for r in list(first.runs)[1:]:
        r._r.getparent().remove(r._r)


def set_bullets_in_frame(tf, lines: list[str]) -> None:
    """Fill a text frame with multiple lines, keeping first paragraph formatting for each."""
    if not lines:
        return
    paras = list(tf.paragraphs)
    template_p = paras[0]
    # collect template formatting from first run
    template_run = template_p.runs[0] if template_p.runs else None
    # clear existing paragraphs
    for p in paras:
        p._p.getparent().remove(p._p)
    # build new paragraphs
    from copy import deepcopy
    for i, line in enumerate(lines):
        new_p = deepcopy(template_p._p)
        # clear children except namespaces
        for child in list(new_p):
            tag = etree.QName(child).localname
            if tag in ("r", "br"):
                new_p.remove(child)
        # add a run
        if template_run is not None:
            new_r = deepcopy(template_run._r)
            # set text
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            t_el = new_r.find(f"{{{ns_a}}}t")
            if t_el is None:
                t_el = etree.SubElement(new_r, f"{{{ns_a}}}t")
            t_el.text = line
            new_p.append(new_r)
        else:
            ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            r = etree.SubElement(new_p, f"{{{ns_a}}}r")
            t = etree.SubElement(r, f"{{{ns_a}}}t")
            t.text = line
        tf._txBody.append(new_p)


def find_title_shape(slide):
    # Prefer the placeholder marked as title
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type is not None and "TITLE" in str(ph_type):
                return shape
    # Fallback: largest text frame near the top
    candidates = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    if not candidates:
        return None
    candidates.sort(key=lambda s: (s.top if s.top is not None else 0))
    return candidates[0]


def find_body_shapes(slide, exclude_shape):
    bodies = []
    for s in slide.shapes:
        if s is exclude_shape:
            continue
        if not s.has_text_frame:
            continue
        # skip footer-y bits
        txt = s.text_frame.text.strip().lower()
        if "maverx.nl" in txt:
            continue
        if len(txt) < 2:
            continue
        bodies.append(s)
    return bodies


def set_notes(slide, notes: dict) -> None:
    """Write the 5-field speaker notes block."""
    block = (
        f"**Aim:** {notes['aim']}\n"
        f"**Time:** {int(notes['time_min'])} min\n"
        f"**Instructions:**\n"
        + "\n".join(f"{i + 1}. {step}" for i, step in enumerate(notes["instructions"]))
        + f"\n**Reflective question:** {notes['reflective_question']}\n"
        f"**Debrief:** {notes['debrief']}\n"
    )
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = block


def build_deck(session_plan: dict, master_path: Path, catalog: dict, out_path: Path) -> None:
    out_prs = Presentation(str(master_path))
    src_slides = list(out_prs.slides)

    # Start from the master file so cloned slides, layouts, media, and themes all
    # belong to the same package. Mixing slide/layout objects from another
    # Presentation creates duplicate package entries on save.
    # remove all existing slides from output (keep masters/layouts/theme)
    sldIdLst = out_prs.slides._sldIdLst
    rid_to_drop = []
    for sldId in list(sldIdLst):
        rid_to_drop.append(sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
        sldIdLst.remove(sldId)
    for rid in rid_to_drop:
        try:
            out_prs.part.drop_rel(rid)
        except Exception:
            pass

    for slide_spec in session_plan["slides"]:
        role = slide_spec.get("role", "default_text")
        try:
            tmpl_idx = resolve_template_idx(catalog, role)
        except Exception as e:
            print(f"[build_deck] unknown role {role!r}: {e}; using default_text", file=sys.stderr)
            tmpl_idx = resolve_template_idx(catalog, "default_text")
        if tmpl_idx >= len(src_slides):
            print(f"[build_deck] tmpl_idx {tmpl_idx} OOR; using 0", file=sys.stderr)
            tmpl_idx = 0
        src = src_slides[tmpl_idx]
        new_slide = clone_slide(out_prs, src)

        # Replace title
        title_shape = find_title_shape(new_slide)
        if title_shape is not None and slide_spec.get("title"):
            replace_text_in_frame(title_shape.text_frame, slide_spec["title"])

        # Replace body: use first non-title text frame; concatenate bullets
        body_shapes = find_body_shapes(new_slide, title_shape)
        body_lines = list(slide_spec.get("body", []))
        if slide_spec.get("subtitle"):
            body_lines = [slide_spec["subtitle"]] + body_lines
        if body_shapes and body_lines:
            set_bullets_in_frame(body_shapes[0].text_frame, body_lines)
            # blank out any remaining body shapes that aren't footers
            for extra in body_shapes[1:]:
                if "maverx" not in extra.text_frame.text.lower():
                    # leave alone — these may be decorative numbered placeholders
                    pass

        # Table (optional)
        if slide_spec.get("table"):
            tbl_spec = slide_spec["table"]
            # Find an existing table in the cloned slide, else skip (don't draw)
            from pptx.shapes.graphfrm import GraphicFrame
            for shape in new_slide.shapes:
                if isinstance(shape, GraphicFrame) and shape.has_table:
                    fill_table(shape.table, tbl_spec)
                    break

        # Notes
        if slide_spec.get("notes"):
            set_notes(new_slide, slide_spec["notes"])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_prs.save(str(out_path))


def fill_table(table, spec: dict) -> None:
    """Best-effort: write headers + rows into an existing table, growing/shrinking rows."""
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    # Headers
    if headers and len(table.rows) >= 1:
        for ci, h in enumerate(headers[: len(table.columns)]):
            cell = table.cell(0, ci)
            cell.text = h
    # Rows
    for ri, row in enumerate(rows):
        target_ri = ri + 1
        if target_ri >= len(table.rows):
            break  # cannot easily add rows; truncate
        for ci, val in enumerate(row[: len(table.columns)]):
            cell = table.cell(target_ri, ci)
            cell.text = str(val)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--session-plan", required=True)
    ap.add_argument("--master", required=True)
    ap.add_argument("--catalog", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    plan = json.loads(Path(args.session_plan).read_text())
    catalog = json.loads(Path(args.catalog).read_text())
    build_deck(plan, Path(args.master), catalog, Path(args.out))
    print(f"wrote {args.out}  ({len(plan['slides'])} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
