#!/usr/bin/env python3
"""QA a session deck + plan: notes completeness, timing, footer, fonts/colors, narrative handshake.

Usage:
    python qa_deck.py --session-plan plan.json --deck deck.pptx [--track track.json --session N] --out qa_report.md
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# Palette from references/house_style.md
APPROVED_HEX = {
    "0D006A", "1A0040", "3F0576", "EF4453", "F48A28", "00B0F0", "262626",
    "BCB3FF", "F7B8C0", "FAD0A8", "9FE6FF", "F2F2F2", "EDE9FF",
    "FEF0F1", "FDEBDB", "E7F9FF", "FFFFFF", "000000",
}
APPROVED_FONTS = {"Space Grotesk", "Raleway", "Calibri", "Calibri Light", "Arial", "+mj-lt", "+mn-lt"}

NOTES_REQUIRED = [
    r"\*\*Aim:\*\*",
    r"\*\*Time:\*\*",
    r"\*\*Instructions:\*\*",
    r"\*\*Reflective question:\*\*",
    r"\*\*Debrief:\*\*",
]


def check_deck(deck_path: Path, plan: dict) -> list[str]:
    prs = Presentation(str(deck_path))
    errs: list[str] = []
    slides = list(prs.slides)
    spec_slides = plan.get("slides", [])

    if len(slides) != len(spec_slides):
        errs.append(f"slide count mismatch: deck={len(slides)} plan={len(spec_slides)}")

    used_fonts: set[str] = set()
    used_hex: set[str] = set()
    notes_time_total = 0
    for i, slide in enumerate(slides):
        # Footer check
        all_text = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        if "maverx.nl" not in all_text.lower():
            errs.append(f"slide {i + 1}: missing 'maverx.nl' footer")
        # Editable text check
        text_frames = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
        if not text_frames:
            errs.append(f"slide {i + 1}: no editable text frame")
        # Notes check
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        for pat in NOTES_REQUIRED:
            if not re.search(pat, notes):
                errs.append(f"slide {i + 1}: notes missing {pat}")
        m = re.search(r"\*\*Time:\*\*\s*(\d+)", notes)
        if m:
            notes_time_total += int(m.group(1))
        # Fonts / colors
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        used_fonts.add(run.font.name)
                    try:
                        if run.font.color and run.font.color.rgb:
                            used_hex.add(str(run.font.color.rgb).upper())
                    except Exception:
                        pass

    # Font check
    bad_fonts = [f for f in used_fonts if not any(af.lower() in f.lower() for af in APPROVED_FONTS)]
    if bad_fonts:
        errs.append(f"non-approved fonts in use: {sorted(bad_fonts)}")

    # Color check
    bad_hex = [h for h in used_hex if h not in APPROVED_HEX]
    if bad_hex:
        errs.append(f"non-approved colors in use (not fatal): {sorted(bad_hex)}")

    # Timing sum check
    target = plan.get("duration_min", 0)
    if target and abs(notes_time_total - target) > 5:
        errs.append(f"sum(notes.time)={notes_time_total} min vs duration_min={target} (tol 5)")

    # Didactic arc check
    blocks = [s.get("block") for s in spec_slides]
    required_order = ["kickoff", "theory", "example", "exercise", "wrapup"]
    first_idx = {b: (blocks.index(b) if b in blocks else -1) for b in required_order}
    missing = [b for b, idx in first_idx.items() if idx == -1]
    if missing:
        errs.append(f"didactic arc missing blocks: {missing}")
    else:
        order = [first_idx[b] for b in required_order]
        if order != sorted(order):
            errs.append(f"didactic arc out of order: {dict(zip(required_order, order))}")

    # Theory ↔ Exercise linkage
    theory_targets = [s.get("exercise_target", "").lower() for s in spec_slides if s.get("role", "").startswith("theory_")]
    exercise_body = " ".join(" ".join(s.get("body", [])) for s in spec_slides if s.get("role", "").startswith("exercise_")).lower()
    for t in theory_targets:
        if t and t not in exercise_body:
            errs.append(f"exercise_target {t!r} not referenced in any exercise slide body")

    return errs


def check_track_handshake(track: dict) -> list[str]:
    errs = []
    sessions = track.get("sessions", [])
    for i in range(1, len(sessions)):
        prev = sessions[i - 1].get("post_bite_artefact")
        curr = sessions[i].get("next_session_pre_bite_expects_from_prior")
        if prev != curr:
            errs.append(f"handshake break: session {sessions[i]['n']} expects {curr!r}, prior produces {prev!r}")
    return errs


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--session-plan", required=True)
    ap.add_argument("--deck", required=True)
    ap.add_argument("--track", default=None)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    plan = json.loads(Path(args.session_plan).read_text())
    errs = check_deck(Path(args.deck), plan)
    if args.track:
        track = json.loads(Path(args.track).read_text())
        errs += check_track_handshake(track)

    report = ["# QA report", "", f"Deck: `{args.deck}`", ""]
    if not errs:
        report.append("All checks passed.")
    else:
        report.append(f"## {len(errs)} issue(s) found")
        for e in errs:
            report.append(f"- {e}")
    Path(args.out).write_text("\n".join(report) + "\n")
    print(f"wrote {args.out}  ({len(errs)} issues)")
    return 1 if errs else 0


if __name__ == "__main__":
    sys.exit(main())
