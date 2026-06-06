#!/usr/bin/env python3
"""Enumerate a Maverx master .pptx so the builder can target layouts by index.

Writes <out>/master_index.json with:
  - layouts: [{idx, name, placeholders: [{idx, name, type}]}]
  - theme_colors: theme color slot names (best-effort)
  - master_fonts: major/minor font names from the theme
  - sample_slides: short text preview of each slide in the master file

Run once after dropping a new master into assets/.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def inspect(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))
    data: dict = {
        "source": pptx_path.name,
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slide_width_in": round(Emu(prs.slide_width).inches, 3),
        "slide_height_in": round(Emu(prs.slide_height).inches, 3),
        "layouts": [],
        "sample_slides": [],
    }

    for idx, layout in enumerate(prs.slide_layouts):
        phs = []
        for ph in layout.placeholders:
            phs.append(
                {
                    "idx": ph.placeholder_format.idx,
                    "name": ph.name,
                    "type": str(ph.placeholder_format.type),
                }
            )
        data["layouts"].append(
            {"idx": idx, "name": layout.name, "placeholders": phs}
        )

    # theme fonts (best effort from XML)
    try:
        theme = prs.slide_master.element.getroottree()
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        # The actual theme lives on the master part; try to dig it up.
        master_part = prs.slide_master.part
        theme_part = None
        for rel in master_part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is not None:
            from lxml import etree

            root = etree.fromstring(theme_part.blob)
            major = root.find(".//a:fontScheme/a:majorFont/a:latin", ns)
            minor = root.find(".//a:fontScheme/a:minorFont/a:latin", ns)
            data["master_fonts"] = {
                "major": major.get("typeface") if major is not None else None,
                "minor": minor.get("typeface") if minor is not None else None,
            }
            clr_scheme = root.find(".//a:clrScheme", ns)
            if clr_scheme is not None:
                data["theme_colors"] = [c.tag.split("}", 1)[-1] for c in clr_scheme]
    except Exception as e:  # noqa: BLE001
        data["theme_introspection_error"] = str(e)

    # sample slides (titles + first text bits) to help the LLM pick layouts
    for i, slide in enumerate(prs.slides):
        title = None
        bits = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt:
                continue
            if title is None and shape == slide.shapes.title if slide.shapes.title else False:
                title = txt
            else:
                bits.append(txt[:120])
        data["sample_slides"].append(
            {
                "slide_idx": i,
                "layout_name": slide.slide_layout.name,
                "title": title or (bits[0] if bits else ""),
                "preview": " | ".join(bits[:4]),
            }
        )

    return data


def main() -> int:
    if len(sys.argv) < 2:
        print("usage: inspect_master.py <master.pptx> [out_dir]", file=sys.stderr)
        return 2
    pptx_path = Path(sys.argv[1])
    out_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else pptx_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    data = inspect(pptx_path)
    out_path = out_dir / "master_index.json"
    out_path.write_text(json.dumps(data, indent=2))
    print(f"wrote {out_path}  ({len(data['layouts'])} layouts, {len(data['sample_slides'])} sample slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
