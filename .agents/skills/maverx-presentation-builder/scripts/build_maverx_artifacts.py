#!/usr/bin/env python3
"""Build Maverx PPTX and DOCX artifacts from a lesson_plan.json file."""

from __future__ import annotations

import argparse
import copy
import json
import re
import shutil
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


PRIMARY_DARK = RGBColor(0x0D, 0x00, 0x6A)
OFF_WHITE = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY = RGBColor(0x26, 0x26, 0x26)


TEMPLATE_SOURCE = {
    "01-deck-title": 1,
    "02-process-slide": 2,
    "03-unstructured-three-section-slide": 3,
    "04-agenda": 4,
    "05-text-slide": 5,
    "06-dark-text-slide": 6,
    "07-complex-layout-slide-1": 7,
    "08-complex-layout-slide-2": 8,
    "09-complex-layout-slide-3": 9,
    "10-longer-text-slide": 10,
    "11-tabular-slide": 11,
    "12-itemized-text-boxes": 12,
    "13-four-section-slide": 13,
    "14-dark-section-title-slide": 14,
    "15-extra-process-timetable-slide": 15,
    "16-three-section-slide": 16,
    "17-theory-topic-slides": 17,
    "18-section-title": 18,
    "19-timeline-process-slide": 19,
    "20-hand-out-slide": 20,
    "21-big-question": 21,
    "22-break-time": 22,
    "23-debrief": 23,
}

SPARSE_TEMPLATE_LIMITS = {
    "14-dark-section-title-slide": 1,
    "17-theory-topic-slides": 1,
    "18-section-title": 1,
    "21-big-question": 1,
    "23-debrief": 1,
}


CONTENT_CANDIDATES = {
    "kick-off": [
        "12-itemized-text-boxes",
        "13-four-section-slide",
        "03-unstructured-three-section-slide",
        "05-text-slide",
        "21-big-question",
    ],
    "theory": [
        "17-theory-topic-slides",
        "19-timeline-process-slide",
        "12-itemized-text-boxes",
        "13-four-section-slide",
        "03-unstructured-three-section-slide",
        "16-three-section-slide",
        "05-text-slide",
    ],
    "example": [
        "12-itemized-text-boxes",
        "13-four-section-slide",
        "03-unstructured-three-section-slide",
        "16-three-section-slide",
        "05-text-slide",
        "19-timeline-process-slide",
    ],
    "exercise": [
        "10-longer-text-slide",
        "13-four-section-slide",
        "12-itemized-text-boxes",
        "03-unstructured-three-section-slide",
        "05-text-slide",
        "16-three-section-slide",
    ],
    "wrap-up": [
        "23-debrief",
        "12-itemized-text-boxes",
        "13-four-section-slide",
        "03-unstructured-three-section-slide",
        "05-text-slide",
        "14-dark-section-title-slide",
    ],
}


@dataclass(frozen=True)
class PlannedSlide:
    template_id: str
    source_slide_number: int
    item: dict[str, Any]


def slugify(value: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", value.lower()).strip("-")
    return slug or "maverx-training"


def first_text_shapes(slide) -> list[Any]:
    return [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text_frame is not None
    ]


def iter_text_shapes(container):
    for shape in container.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            yield shape
        if hasattr(shape, "shapes"):
            yield from iter_text_shapes(shape)


def clear_text(shape) -> None:
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        shape.text_frame.clear()


def clear_unused_text_shapes(slide, keep: set[int]) -> None:
    top_level_shapes = first_text_shapes(slide)
    for index, shape in enumerate(top_level_shapes):
        if index not in keep:
            clear_text(shape)
    top_level_elements = {shape.element for shape in top_level_shapes}
    for shape in iter_text_shapes(slide):
        if shape.element not in top_level_elements:
            clear_text(shape)


def set_text(
    shape,
    paragraphs: str | list[str],
    *,
    font: str = "Raleway",
    size: float = 15,
    color: RGBColor = PRIMARY_DARK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    if isinstance(paragraphs, str):
        paragraphs = [paragraphs]
    text_frame = shape.text_frame
    text_frame.clear()
    for index, line in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.font.name = font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def set_paragraphs(
    shape,
    paragraphs: list[dict[str, Any]],
    *,
    default_font: str = "Raleway",
    default_size: float = 15,
    default_color: RGBColor = PRIMARY_DARK,
    default_bold: bool = False,
    default_align=PP_ALIGN.LEFT,
) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    for index, spec in enumerate(paragraphs):
        paragraph = (
            text_frame.paragraphs[0]
            if index == 0
            else text_frame.add_paragraph()
        )
        paragraph.text = spec.get("text", "")
        paragraph.alignment = spec.get("align", default_align)
        font_name = spec.get("font", default_font)
        font_size = spec.get("size", default_size)
        font_color = spec.get("color", default_color)
        font_bold = spec.get("bold", default_bold)
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = font_bold
        paragraph.font.color.rgb = font_color
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = font_bold
            run.font.color.rgb = font_color


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def unique_lines(values: list[str]) -> list[str]:
    seen: set[str] = set()
    lines: list[str] = []
    for value in values:
        line = normalize_text(value)
        if not line:
            continue
        key = line.lower()
        if key in seen:
            continue
        seen.add(key)
        lines.append(line)
    return lines


def fallback_points(item: dict[str, Any]) -> list[str]:
    didactic_block = item.get("didactic_block")
    defaults = {
        "kick-off": [
            "Name the task clearly.",
            "State what good output should help you decide or do.",
            "Surface the assumption that matters most.",
        ],
        "theory": [
            "Define the idea in plain language.",
            "Show the mechanism that makes it useful.",
            "State the practical takeaway for consultants.",
        ],
        "example": [
            "Show the original situation.",
            "Explain the move that improves it.",
            "Name the difference in the result.",
        ],
        "exercise": [
            "Set the task participants need to complete.",
            "State what to inspect or revise.",
            "Make the review check explicit.",
        ],
        "wrap-up": [
            "Name the pattern to remember.",
            "State the takeaway in work language.",
            "Choose the next action after the session.",
        ],
    }
    return defaults.get(didactic_block, [])


def content_pool(item: dict[str, Any], limit: int = 6) -> list[str]:
    candidates = list(item.get("suggested_content", []))
    candidates.extend(
        [
            item.get("learning_purpose", ""),
            item.get("key_message", ""),
        ]
    )
    candidates.extend(fallback_points(item))
    return unique_lines(candidates)[:limit]


def lead_statement(item: dict[str, Any]) -> str:
    for candidate in [
        item.get("key_message", ""),
        item.get("learning_purpose", ""),
        item.get("title", ""),
    ]:
        line = normalize_text(candidate)
        if line:
            return line
    return "Use the slide to move the session forward."


def short_label(value: str, max_words: int = 4, max_chars: int = 28) -> str:
    words = normalize_text(value).split()
    label = " ".join(words[:max_words]).strip(":,. ")
    return label[:max_chars] or "Key point"


def set_notes(slide, item: dict[str, Any], session: dict[str, Any]) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    title = item.get("title") or item.get("message") or session["title"]
    reliability = item.get("reliability", {})
    notes.text = (
        f"Aim: {item.get('learning_purpose', 'Support the session flow.')} "
        f"Slide focus: {title}\n\n"
        f"Time: Use the timing from the lesson plan for the "
        f"{item.get('didactic_block', item.get('slide_type', 'slide'))} block.\n\n"
        f"Instructions: Present the key message, connect it to the session "
        f"purpose, and keep the discussion anchored in participant work.\n\n"
        f"Key discussion points: {format_points(item)}\n\n"
        f"Link to reality: Relate this slide to {session['brief_for_trainer']['session_purpose']}\n\n"
        f"Debrief & Summary: {item.get('key_message', title)}\n\n"
        f"Reliability: {reliability.get('score', 'n/a')} "
        f"({reliability.get('review_priority', 'n/a')} review). "
        f"{reliability.get('rationale', '')}"
    )


def format_points(item: dict[str, Any]) -> str:
    if item.get("slide_type") == "content":
        points = item.get("suggested_content", [])
        if points:
            return "; ".join(points[:4])
        return item.get("key_message", "")
    if item.get("slide_type") == "break":
        return f"Break slide; duration {item.get('duration_min')} minutes"
    if item.get("slide_type") == "handout":
        return (
            f"Handout work slide; time budget "
            f"{item.get('time_budget_min')} minutes"
        )
    return ""


def choose_templates(deck_outline: list[dict[str, Any]]) -> list[PlannedSlide]:
    block_usage_count: dict[str, int] = {}
    template_usage_count: dict[str, int] = {}
    planned: list[PlannedSlide] = []
    for item in deck_outline:
        slide_type = item.get("slide_type", "content")
        if item.get("slide_n") == 1:
            template_id = "01-deck-title"
        elif slide_type == "break":
            template_id = "22-break-time"
        elif slide_type == "handout":
            template_id = "20-hand-out-slide"
        elif "debrief" in item.get("title", "").lower():
            template_id = "23-debrief"
        else:
            didactic_block = item.get("didactic_block", "")
            candidates = CONTENT_CANDIDATES.get(
                didactic_block, ["05-text-slide"]
            )
            start_index = block_usage_count.get(didactic_block, 0)
            template_id = candidates[start_index % len(candidates)]
            for offset in range(len(candidates)):
                candidate = candidates[(start_index + offset) % len(candidates)]
                limit = SPARSE_TEMPLATE_LIMITS.get(candidate)
                if limit is not None and template_usage_count.get(candidate, 0) >= limit:
                    continue
                template_id = candidate
                break
            block_usage_count[didactic_block] = (
                block_usage_count.get(didactic_block, 0) + 1
            )
        template_usage_count[template_id] = template_usage_count.get(template_id, 0) + 1
        planned.append(
            PlannedSlide(
                template_id=template_id,
                source_slide_number=TEMPLATE_SOURCE[template_id],
                item=item,
            )
        )
    return planned


def replace_relationship_ids(element, rel_map: dict[str, str]) -> None:
    for node in element.iter():
        for attr_name, attr_value in list(node.attrib.items()):
            if attr_value in rel_map:
                node.attrib[attr_name] = rel_map[attr_value]


def clone_slide(prs: Presentation, source_slide) -> Any:
    clone = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(clone.shapes):
        shape.element.getparent().remove(shape.element)

    rel_map: dict[str, str] = {}
    for rel in source_slide.part.rels.values():
        if rel.reltype.endswith("/slideLayout") or rel.reltype.endswith("/notesSlide"):
            continue
        rel_map[rel.rId] = clone.part.rels._add_relationship(
            rel.reltype,
            rel._target,
            rel.is_external,
        )

    for shape in source_slide.shapes:
        element = copy.deepcopy(shape.element)
        replace_relationship_ids(element, rel_map)
        clone.shapes._spTree.insert_element_before(element, "p:extLst")
    return clone


def remove_slides_except(prs: Presentation, slides_to_keep: list[Any]) -> None:
    all_sld_ids = list(prs.slides._sldIdLst)
    keep_slide_ids = {slide.slide_id for slide in slides_to_keep}
    for sld_id in all_sld_ids:
        if int(sld_id.id) not in keep_slide_ids:
            prs.slides._sldIdLst.remove(sld_id)
            prs.part.drop_rel(sld_id.rId)


def fill_slide(slide, planned: PlannedSlide, session: dict[str, Any]) -> None:
    item = planned.item
    template_id = planned.template_id
    slide_type = item.get("slide_type", "content")
    if slide_type == "break":
        fill_break(slide, item)
    elif slide_type == "handout":
        fill_handout(slide, item)
    elif template_id == "01-deck-title":
        fill_title(slide, item, session)
    elif template_id == "04-agenda":
        fill_agenda(slide, session)
    elif template_id == "03-unstructured-three-section-slide":
        fill_unstructured_three_section(slide, item)
    elif template_id == "05-text-slide":
        fill_text_slide(slide, item)
    elif template_id == "10-longer-text-slide":
        fill_long_text(slide, item)
    elif template_id == "12-itemized-text-boxes":
        fill_itemized_text_boxes(slide, item)
    elif template_id == "13-four-section-slide":
        fill_four_section(slide, item)
    elif template_id == "17-theory-topic-slides":
        fill_theory_opener(slide, item)
    elif template_id == "19-timeline-process-slide":
        fill_timeline(slide, item)
    elif template_id == "16-three-section-slide":
        fill_three_section(slide, item)
    elif template_id == "23-debrief":
        fill_debrief(slide, item)
    else:
        fill_generic_content(slide, item, template_id)
    set_notes(slide, item, session)


def fill_title(slide, item: dict[str, Any], session: dict[str, Any]) -> None:
    shapes = first_text_shapes(slide)
    set_text(shapes[0], session["title"], font="Space Grotesk", size=36, color=OFF_WHITE, bold=True)
    subtitle = item.get("key_message") or session["brief_for_trainer"]["session_purpose"]
    set_text(shapes[-1], subtitle[:115], font="Space Grotesk", size=22, color=OFF_WHITE)


def fill_agenda(slide, session: dict[str, Any]) -> None:
    shapes = first_text_shapes(slide)
    set_text(shapes[0], "Agenda", font="Space Grotesk", size=33, color=PRIMARY_DARK, bold=True)
    blocks = session["didactic_arc"]
    text_frame = shapes[1].text_frame
    text_frame.clear()
    text_frame.paragraphs[0].text = ""
    for block in blocks:
        paragraph = text_frame.add_paragraph()
        paragraph.text = block["block"].title()
        paragraph.font.name = "Space Grotesk"
        paragraph.font.size = Pt(18.5)
        paragraph.font.color.rgb = PRIMARY_DARK
        paragraph = text_frame.add_paragraph()
        paragraph.text = block["learning_purpose"][:82]
        paragraph.font.name = "Raleway"
        paragraph.font.size = Pt(9.5)
        paragraph.font.color.rgb = PRIMARY_DARK


def fill_unstructured_three_section(slide, item: dict[str, Any]) -> None:
    shapes = first_text_shapes(slide)
    set_text(
        shapes[0],
        item["title"],
        font="Space Grotesk",
        size=33,
        color=PRIMARY_DARK,
        bold=True,
    )
    pool = content_pool(item, limit=6)
    sections = [
        ("Core idea", [lead_statement(item), *pool[:1]]),
        ("Why it matters", pool[1:3] or [item.get("learning_purpose", "")]),
        ("Use or check", pool[3:6] or fallback_points(item)[:2]),
    ]
    paragraphs: list[dict[str, Any]] = []
    for title, lines in sections:
        if paragraphs:
            paragraphs.append({"text": "", "size": 7})
        paragraphs.append(
            {
                "text": title,
                "font": "Raleway",
                "size": 19,
                "bold": True,
                "color": PRIMARY_DARK,
            }
        )
        for line in unique_lines(lines)[:3]:
            paragraphs.append({"text": f"- {line}", "size": 12.5})
    set_paragraphs(shapes[1], paragraphs, default_color=PRIMARY_DARK)


def fill_text_slide(slide, item: dict[str, Any]) -> None:
    shapes = first_text_shapes(slide)
    set_text(
        shapes[0],
        item["title"],
        font="Space Grotesk",
        size=33,
        color=PRIMARY_DARK,
        bold=True,
    )
    lead = lead_statement(item)
    points = [
        point for point in content_pool(item, limit=5)
        if normalize_text(point).lower() != normalize_text(lead).lower()
    ][:4]
    paragraphs = [{"text": lead, "size": 17, "bold": True}]
    if points:
        paragraphs.append({"text": "", "size": 7})
    for index, point in enumerate(points, start=1):
        paragraphs.append({"text": f"{index}. {point}", "size": 14.5})
    set_paragraphs(shapes[1], paragraphs, default_color=PRIMARY_DARK)


def fill_theory_opener(slide, item: dict[str, Any]) -> None:
    set_text(
        first_text_shapes(slide)[0],
        item["title"],
        font="Space Grotesk",
        size=32,
        color=PRIMARY_DARK,
        bold=True,
    )


def fill_long_text(slide, item: dict[str, Any]) -> None:
    shapes = first_text_shapes(slide)
    set_text(
        shapes[0],
        item["title"],
        font="Space Grotesk",
        size=33,
        color=PRIMARY_DARK,
        bold=True,
    )
    points = content_pool(item, limit=6)
    paragraphs = [{"text": lead_statement(item), "size": 15.5, "bold": True}]
    for index, point in enumerate(points[:5], start=1):
        paragraphs.append({"text": f"{index}. {point}", "size": 12.5})
    set_paragraphs(shapes[1], paragraphs, default_color=PRIMARY_DARK)


def fill_timeline(slide, item: dict[str, Any]) -> None:
    shapes = first_text_shapes(slide)
    set_text(
        shapes[0],
        item["title"],
        font="Space Grotesk",
        size=33,
        color=PRIMARY_DARK,
        bold=True,
    )
    suggested = item.get("suggested_content", [])
    labels = (suggested + ["Step 1", "Step 2", "Step 3", "Step 4", "Step 5"])[:5]
    # Template 19 uses text shapes 1-5 for numbered markers and 6-10 for
    # labels beneath the markers.
    for marker, number in zip(shapes[1:6], ["1", "2", "3", "4", "5"]):
        set_text(
            marker,
            number,
            font="Space Grotesk",
            size=24,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
    for label_shape, label in zip(shapes[6:11], labels):
        set_text(
            label_shape,
            shorten(label, 11),
            font="Space Grotesk",
            size=7.5,
            color=DARK_GREY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    if len(shapes) >= 9:
        set_text(shapes[-3], ["DEFINITION", item["key_message"]], font="Raleway", size=10.5, color=OFF_WHITE)
        set_text(shapes[-2], "Main takeaway", font="Space Grotesk", size=15, color=PRIMARY_DARK, bold=True)
        set_text(shapes[-1], item["learning_purpose"], font="Raleway", size=14, color=PRIMARY_DARK)


def fill_itemized_text_boxes(slide, item: dict[str, Any]) -> None:
    shapes = list(iter_text_shapes(slide))
    set_text(
        shapes[2],
        item["title"],
        font="Space Grotesk",
        size=31,
        color=PRIMARY_DARK,
        bold=True,
    )
    card_shapes = sorted(
        [shapes[5], shapes[17], shapes[18], shapes[19], shapes[20]],
        key=lambda shape: (shape.top, shape.left),
    )
    number_shapes = sorted(
        [shapes[12], shapes[13], shapes[14], shapes[15], shapes[16]],
        key=lambda shape: (shape.top, shape.left),
    )
    points = content_pool(item, limit=5)
    while len(points) < 5:
        points.append(fallback_points(item)[len(points) % len(fallback_points(item))])
    for index, shape in enumerate(number_shapes, start=1):
        set_text(
            shape,
            str(index),
            font="Space Grotesk",
            size=16,
            color=PRIMARY_DARK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    for shape, point in zip(card_shapes, points):
        set_paragraphs(
            shape,
            [
                {"text": short_label(point).upper(), "size": 11, "bold": True},
                {"text": point, "size": 11.5},
            ],
            default_color=PRIMARY_DARK,
        )


def fill_three_section(slide, item: dict[str, Any]) -> None:
    shapes = first_text_shapes(slide)
    if len(shapes) < 9:
        fill_generic_content(slide, item, "16-three-section-slide")
        return

    set_text(
        shapes[0],
        item["title"],
        font="Space Grotesk",
        size=31,
        color=PRIMARY_DARK,
        bold=True,
    )
    set_text(
        shapes[1],
        item["key_message"],
        font="Raleway",
        size=11.5,
        color=PRIMARY_DARK,
    )

    points = item.get("suggested_content", [])[:3]
    while len(points) < 3:
        points.append("")
    card_groups = [
        (2, 3, 4, WHITE),
        (5, 6, 7, PRIMARY_DARK),
        (8, 9, 10, WHITE),
    ]
    for point, (background, line_1, line_2, color) in zip(points, card_groups):
        set_text(
            shapes[background],
            point,
            font="Raleway",
            size=17,
            color=color,
        )
        for index in [line_1, line_2]:
            if index < len(shapes):
                clear_text(shapes[index])


def fill_four_section(slide, item: dict[str, Any]) -> None:
    shapes = list(iter_text_shapes(slide))
    set_text(
        shapes[27],
        item["title"],
        font="Space Grotesk",
        size=31,
        color=PRIMARY_DARK,
        bold=True,
    )
    sections = [
        ("Core idea", lead_statement(item)),
        ("Why it matters", item.get("learning_purpose", "")),
        ("In practice", " | ".join(content_pool(item, limit=2))),
        ("Check or next step", " | ".join(content_pool(item, limit=5)[2:5])),
    ]
    body_defaults = fallback_points(item)
    if not sections[3][1]:
        sections[3] = ("Check or next step", " | ".join(body_defaults[:2]))
    header_shapes = [shapes[7], shapes[17], shapes[3], shapes[24]]
    body_shapes = [shapes[13], shapes[19], shapes[20], shapes[26]]
    for (header, body), header_shape, body_shape in zip(
        sections, header_shapes, body_shapes
    ):
        set_text(
            header_shape,
            header,
            font="Space Grotesk SemiBold",
            size=14,
            color=PRIMARY_DARK,
            bold=True,
        )
        lines = unique_lines(body.split("|"))[:3]
        set_paragraphs(
            body_shape,
            [{"text": f"- {line}" if len(lines) > 1 else line, "size": 11.5} for line in lines],
            default_color=PRIMARY_DARK,
        )


def fill_break(slide, item: dict[str, Any]) -> None:
    message = item.get("message", "Break time!")
    duration = item.get("duration_min")
    if re.search(r"\b\d+\s*minutes?\b", message, flags=re.IGNORECASE):
        message = "Break time!"
    set_text(
        first_text_shapes(slide)[0],
        [message, f"{duration} minutes"],
        font="Space Grotesk",
        size=33,
        color=PRIMARY_DARK,
        bold=True,
    )
    paragraphs = first_text_shapes(slide)[0].text_frame.paragraphs
    if len(paragraphs) > 1:
        paragraphs[1].font.size = Pt(18)
        paragraphs[1].font.bold = False


def fill_handout(slide, item: dict[str, Any]) -> None:
    shapes = first_text_shapes(slide)
    set_text(shapes[0], "Handout work", font="Space Grotesk", size=33, color=WHITE, bold=True)
    set_text(
        shapes[1],
        [item.get("message", "Work on the handout."), f"Time budget: {item.get('time_budget_min')} minutes"],
        font="Raleway",
        size=15,
        color=WHITE,
    )


def fill_debrief(slide, item: dict[str, Any]) -> None:
    title = item.get("title", "Debrief")
    size = 29 if len(title) > 34 else 33
    set_text(
        first_text_shapes(slide)[0],
        title,
        font="Space Grotesk",
        size=size,
        color=OFF_WHITE,
        bold=True,
    )


def fill_generic_content(slide, item: dict[str, Any], template_id: str) -> None:
    shapes = first_text_shapes(slide)
    if not shapes:
        return
    dark = template_id in {"06-dark-text-slide", "14-dark-section-title-slide"}
    title_color = OFF_WHITE if dark else PRIMARY_DARK
    body_color = OFF_WHITE if dark else PRIMARY_DARK
    set_text(shapes[0], item["title"], font="Space Grotesk", size=33, color=title_color, bold=True)
    body = [lead_statement(item), *content_pool(item, limit=4)]
    keep = {0}
    for index, (shape, text) in enumerate(zip(shapes[1:], body), start=1):
        set_text(shape, text, font="Raleway", size=15, color=body_color)
        keep.add(index)
    clear_unused_text_shapes(slide, keep)


def shorten(value: str, max_len: int) -> str:
    value = re.sub(r"\s+", " ", value).strip()
    if len(value) <= max_len:
        return value
    words = value.split()
    return words[0][:max_len]


def build_session_deck(
    lesson_plan: dict[str, Any],
    session: dict[str, Any],
    template_path: Path,
    out_dir: Path,
) -> Path:
    planned = choose_templates(session["deck_outline"])
    prs = Presentation(template_path)
    source_slides = list(prs.slides)
    cloned_slides = [
        clone_slide(prs, source_slides[slide.source_slide_number - 1])
        for slide in planned
    ]
    remove_slides_except(prs, cloned_slides)
    for slide, planned_slide in zip(cloned_slides, planned):
        fill_slide(slide, planned_slide, session)
    output = out_dir / (
        f"{lesson_plan['training']['slug']}-session-"
        f"{session['session_n']:02d}.pptx"
    )
    prs.save(output)
    Presentation(output)
    return output


def docx_escape_lines(lines: list[str]) -> str:
    paragraphs = []
    for line in lines:
        safe = escape(line)
        paragraphs.append(f"<w:p><w:r><w:t>{safe}</w:t></w:r></w:p>")
    return "".join(paragraphs)


def write_docx(path: Path, title: str, sections: list[tuple[str, list[str]]]) -> None:
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
        "</Relationships>"
    )
    body_lines = [title, ""]
    for heading, lines in sections:
        body_lines.extend([heading, *lines, ""])
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body>{docx_escape_lines(body_lines)}"
        '<w:sectPr><w:pgSz w:w="11906" w:h="16838"/>'
        '<w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440"/>'
        "</w:sectPr></w:body></w:document>"
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as package:
        package.writestr("[Content_Types].xml", content_types)
        package.writestr("_rels/.rels", rels)
        package.writestr("word/document.xml", document)


def build_docs(lesson_plan: dict[str, Any], session: dict[str, Any], out_dir: Path) -> list[Path]:
    slug = lesson_plan["training"]["slug"]
    session_prefix = f"{slug}-session-{session['session_n']:02d}"
    outputs: list[Path] = []
    pre = session["pre_bite"]
    pre_path = out_dir / f"{session_prefix}-pre-bite.docx"
    write_docx(
        pre_path,
        f"Pre-bite: {session['title']}",
        [
            ("Purpose", [pre["purpose"]]),
            ("Time", [f"{pre['time_min']} minutes"]),
            ("Participant task", [pre["participant_task"]]),
            (
                "Resources",
                [
                    f"{resource['title']} - {resource['url_or_reference']} ({resource['why_it_is_included']})"
                    for resource in pre.get("resources", [])
                ],
            ),
        ],
    )
    outputs.append(pre_path)
    post = session["post_bite"]
    post_path = out_dir / f"{session_prefix}-post-bite.docx"
    write_docx(
        post_path,
        f"Post-bite: {session['title']}",
        [
            ("Purpose", [post["purpose"]]),
            ("Time", [f"{post['time_min']} minutes"]),
            ("Participant task", [post["participant_task"]]),
            (
                "Resources",
                [
                    f"{resource['title']} - {resource['url_or_reference']} ({resource['why_it_is_included']})"
                    for resource in post.get("resources", [])
                ],
            ),
        ],
    )
    outputs.append(post_path)
    handout = session.get("handout")
    if handout:
        handout_path = out_dir / f"{session_prefix}-handout.docx"
        lines = [f"Time budget: {handout['time_budget_min']} minutes", handout["purpose"]]
        for activity in handout.get("activities", []):
            lines.append(activity["activity_title"])
            lines.extend(activity["instructions"])
            lines.append(f"Participant output: {activity['participant_output']}")
        write_docx(handout_path, f"Handout: {session['title']}", [("Activity", lines)])
        outputs.append(handout_path)
    return outputs


def validate_lesson_plan(plan: dict[str, Any]) -> None:
    for key in ["training", "sessions"]:
        if key not in plan:
            raise ValueError(f"lesson_plan.json is missing required key: {key}")
    for session in plan["sessions"]:
        for key in ["session_n", "title", "deck_outline", "pre_bite", "post_bite"]:
            if key not in session:
                raise ValueError(f"session is missing required key: {key}")
        if not session["deck_outline"]:
            raise ValueError(f"session {session['session_n']} has no deck_outline")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("lesson_plan", type=Path)
    parser.add_argument("--template", type=Path, required=True)
    parser.add_argument("--out-dir", type=Path)
    parser.add_argument("--clean", action="store_true")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    plan = json.loads(args.lesson_plan.read_text(encoding="utf-8"))
    validate_lesson_plan(plan)
    out_dir = args.out_dir or args.lesson_plan.parent / "presentation_artifacts"
    if args.clean and out_dir.exists():
        shutil.rmtree(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    manifest: dict[str, Any] = {
        "created_at": datetime.now(timezone.utc).isoformat(),
        "lesson_plan": str(args.lesson_plan),
        "outputs": [],
    }
    for session in plan["sessions"]:
        deck = build_session_deck(plan, session, args.template, out_dir)
        docs = build_docs(plan, session, out_dir)
        manifest["outputs"].append(
            {
                "session_n": session["session_n"],
                "pptx": str(deck),
                "docx": [str(path) for path in docs],
            }
        )
    manifest_path = out_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(manifest_path)


if __name__ == "__main__":
    main()
