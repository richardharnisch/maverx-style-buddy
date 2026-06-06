"""Skill: create a new presentation from a branded template or blank."""

import logging
import uuid

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from src.style_guides.loader import load_style_guide, resolve_template_path
from src.skills import _store

log = logging.getLogger(__name__)

TOOL_SPEC = {
    "type": "function",
    "function": {
        "name": "create_presentation",
        "description": "Create a new PowerPoint presentation. Uses the named style guide's branded template if available.",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Presentation title"},
                "style_guide": {"type": "string", "description": "Name of the style guide to apply (default: 'default')"},
            },
            "required": ["title"],
        },
    },
}


def _clear_slides(prs: Presentation) -> None:
    """Remove all content slides from a template, preserving master and layouts."""
    slide_ids = list(prs.slides._sldIdLst)
    for sId in slide_ids:
        rId = sId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sId)


def execute(title: str, style_guide: str = "default") -> dict:
    log.info("Creating presentation: title='%s' style_guide='%s'", title, style_guide)
    guide = load_style_guide(style_guide)
    template_path = resolve_template_path(guide)

    if template_path is not None:
        log.debug("Opening template: %s", template_path.name)
        prs = Presentation(str(template_path))
        n_slides = len(prs.slides)
        _clear_slides(prs)
        log.debug("Cleared %d existing slides from template", n_slides)
    else:
        log.debug("No template found — building blank presentation")
        prs = Presentation()
        if guide.slide_size:
            from pptx.util import Emu
            prs.slide_width = Emu(guide.slide_size.width_emu)
            prs.slide_height = Emu(guide.slide_size.height_emu)
            log.debug("Slide size set: %dx%d EMU", guide.slide_size.width_emu, guide.slide_size.height_emu)
        if guide.colors:
            bg = prs.slide_master.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(guide.colors.background.lstrip("#"))
            log.debug("Master background set to %s", guide.colors.background)

    pid = str(uuid.uuid4())[:8]
    _store.put(pid, prs, guide, template_path)
    log.info("Presentation created: id=%s template=%s", pid, guide.template_pptx or "blank")

    return {
        "presentation_id": pid,
        "template": guide.template_pptx or "blank",
        "style_guide": guide.name,
    }
