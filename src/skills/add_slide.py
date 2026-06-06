"""Skill: append a slide to an in-progress presentation using a named layout template."""

import logging

from pptx.util import Pt

from src.skills import _store
from src.style_guides.loader import load_slide_templates
from src.style_guides.schema import DecorativeShape, SlideTemplate, TableArea, TextArea

log = logging.getLogger(__name__)

TOOL_SPEC = {
    "type": "function",
    "function": {
        "name": "add_slide",
        "description": (
            "Add a slide to the presentation using a named layout template. "
            "Call list_layouts first to discover valid layout keys and their descriptions. "
            "Slides with image areas will automatically include their default template images. "
            "Use 'image_overrides' to replace a specific image slot with a different image from list_images."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "presentation_id": {"type": "string"},
                "layout": {
                    "type": "string",
                    "description": "Layout key from list_layouts (e.g. 'title_slide_1', 'titeldia_3').",
                },
                "title": {"type": "string", "description": "Text for the title area."},
                "body": {
                    "type": "string",
                    "description": (
                        "Text for the body/subtitle area. "
                        "For two-column layouts, separate left and right with \\n---\\n."
                    ),
                },
                "speaker_notes": {"type": "string"},
                "image_overrides": {
                    "type": "array",
                    "description": (
                        "Optional: override default images. Each item: "
                        "{\"role\": \"background\", \"image_key\": \"img_abc123\"}. "
                        "Get image keys from list_images."
                    ),
                    "items": {
                        "type": "object",
                        "properties": {
                            "role": {"type": "string"},
                            "image_key": {"type": "string"},
                        },
                        "required": ["role", "image_key"],
                    },
                },
            },
            "required": ["presentation_id", "layout"],
        },
    },
}


# ── Layout lookup ──────────────────────────────────────────────────────────────

def _find_pptx_layout(prs, layout_name: str):
    """Search all slide masters for a layout by name."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout
    log.warning("Layout '%s' not found in any master; using first available", layout_name)
    return prs.slide_masters[0].slide_layouts[0]


# ── Text box creation ──────────────────────────────────────────────────────────

# MSO auto shape preset name → integer ID used by add_shape()
_PRESET_ID: dict[str, int] = {
    "rect": 1,
    "roundRect": 5,
    "ellipse": 9,
    "diamond": 4,
    "triangle": 7,
    "snip1Rect": 88,
    "snip2SameRect": 89,
    "snipRoundRect": 84,
}


def _hex_to_rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_text_box(slide, area: TextArea, text: str) -> None:
    from pptx.enum.text import PP_ALIGN

    _ALIGN_MAP = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    if area.shape_preset and area.shape_preset in _PRESET_ID:
        shape_id = _PRESET_ID[area.shape_preset]
        shape_obj = slide.shapes.add_shape(shape_id, area.left, area.top, area.width, area.height)
        if area.corner_radius is not None:
            try:
                shape_obj.adjustments[0] = area.corner_radius
            except Exception:
                pass
        if area.bg_color:
            try:
                shape_obj.fill.solid()
                shape_obj.fill.fore_color.rgb = _hex_to_rgb(area.bg_color)
            except Exception:
                pass
        else:
            try:
                shape_obj.fill.background()  # transparent
            except Exception:
                pass
        tf = shape_obj.text_frame
    else:
        box = slide.shapes.add_textbox(area.left, area.top, area.width, area.height)
        tf = box.text_frame

    tf.word_wrap = True
    tf.text = text

    align_val = _ALIGN_MAP.get(area.alignment) if area.alignment else None
    for para in tf.paragraphs:
        if align_val is not None:
            para.alignment = align_val
        for run in para.runs:
            if area.font_name:
                run.font.name = area.font_name
            if area.font_size_pt:
                run.font.size = Pt(area.font_size_pt)
            if area.bold is not None:
                run.font.bold = area.bold
            if area.italic is not None:
                run.font.italic = area.italic
            if area.underline is not None:
                run.font.underline = area.underline
            if area.color:
                try:
                    run.font.color.rgb = _hex_to_rgb(area.color)
                except Exception:
                    pass


# ── Text distribution and sizing ──────────────────────────────────────────────

def _estimate_word_capacity(area: TextArea) -> int:
    """Estimate how many words would fit in a text box based on dimensions and font size.

    Heuristic: typical proportional font, ~0.55 of font size = char width,
    ~7 chars per word, line height ~1.15x font size.
    """
    if area.width <= 0 or area.height <= 0:
        return 0

    font_size = area.font_size_pt or 12.0

    # Convert box dimensions from EMU to points (914400 EMU = 72 points)
    width_pt = area.width / 914400 * 72
    height_pt = area.height / 914400 * 72

    # Character width heuristic: ~0.55 of font size for proportional fonts
    # (accounts for monospace being wider, but most fonts are narrower)
    approx_char_width_pt = font_size * 0.55

    # Chars per line with ~5% margin for padding
    chars_per_line = max(1, int((width_pt * 0.95) / approx_char_width_pt))

    # Line height: font_size * 1.15 (typical single-line spacing in PowerPoint)
    line_height_pt = font_size * 1.15

    # Available lines with 10% margin for top/bottom padding
    available_lines = max(1, int((height_pt * 0.9) / line_height_pt))

    # Total characters available
    chars_available = chars_per_line * available_lines

    # Average 7 characters per word (English average: 4.7, but we add spaces)
    words_available = max(1, chars_available // 7)

    log.debug(
        "Text capacity: role='%s' → %d words (font %.1fpt, box %.0fx%.0fpt, %d ch/ln, %d lines)",
        area.role, words_available, font_size, width_pt, height_pt,
        chars_per_line, available_lines
    )
    return words_available


def _distribute_text_by_role(template: SlideTemplate, title: str, body: str) -> dict[int, str]:
    """Distribute content across text areas, accounting for capacity of each box.

    Returns dict mapping text_area index to the content for that box.
    Roles with multiple boxes get content distributed across them based on capacity.
    """
    distribution = {}
    role_indices: dict[str, list[int]] = {}

    # Index all text areas by role
    for idx, area in enumerate(template.text_areas):
        if area.role not in role_indices:
            role_indices[area.role] = []
        role_indices[area.role].append(idx)

    # Build the content pool
    content_pool = {
        "title": title,
        "subtitle": body,
        "body": body,
        "left_column": body.partition("\n---\n")[0].strip(),
        "right_column": body.partition("\n---\n")[2].strip(),
        "caption": "",
    }

    # Distribute content for each role
    for role, indices in role_indices.items():
        content = content_pool.get(role, "")
        words = content.split() if content else []
        word_idx = 0

        for area_idx in indices:
            area = template.text_areas[area_idx]
            capacity = _estimate_word_capacity(area)

            if word_idx >= len(words):
                # No more content for this role
                distribution[area_idx] = ""
                continue

            # Take up to capacity words for this box
            box_words = words[word_idx : word_idx + capacity]
            distribution[area_idx] = " ".join(box_words)
            word_idx += len(box_words)

            log.debug(
                "Area[%d] role='%s': %d words (capacity %d)",
                area_idx, role, len(box_words), capacity
            )

        if word_idx < len(words):
            log.warning(
                "Role '%s': %d words did not fit (%.0f%% overflow)",
                role, len(words) - word_idx, 100.0 * (len(words) - word_idx) / len(words)
            )

    return distribution


# ── Role → content mapping ────────────────────────────────────────────────────

def _build_role_map(title: str, body: str) -> dict[str, str]:
    left, _, right = body.partition("\n---\n")
    return {
        "title": title,
        "subtitle": body,
        "body": body,
        "left_column": left.strip(),
        "right_column": right.strip() if right else "",
        "caption": "",  # captions are decorative; skip if no explicit content
    }


# ── Image placement ───────────────────────────────────────────────────────────

def _place_images(slide, template, guide, image_overrides: list[dict]) -> list[str]:
    """Add image shapes to slide from template image areas. Returns list of placed roles."""
    if not template.image_areas:
        return []

    from src.style_guides.loader import PROJECT_ROOT

    layouts_dir = PROJECT_ROOT / guide.layouts_dir if guide.layouts_dir else None
    if not layouts_dir:
        return []
    images_dir = layouts_dir / "images"

    override_map = {o["role"]: o["image_key"] for o in (image_overrides or [])}
    placed = []

    for ia in template.image_areas:
        image_key = override_map.get(ia.role, ia.image_key)
        if not image_key:
            continue
        # Find the image file (any extension)
        matches = list(images_dir.glob(f"{image_key}.*")) if images_dir.exists() else []
        img_path = next((m for m in matches if m.suffix != ".yaml"), None)
        if not img_path:
            log.warning("Image file not found for key '%s'", image_key)
            continue
        try:
            slide.shapes.add_picture(str(img_path), ia.left, ia.top, ia.width, ia.height)
            placed.append(ia.role)
            log.debug("Placed image role='%s' key='%s'", ia.role, image_key)
        except Exception as e:
            log.warning("Could not place image '%s': %s", image_key, e)

    return placed


# ── Decorative shape placement ─────────────────────────────────────────────────

def _place_decorative_shapes(slide, template: SlideTemplate) -> list[str]:
    """Add decorative shapes (rectangles, circles, arrows, lines) to slide."""
    if not template.decorative_shapes:
        return []

    placed = []
    _PRESET_ID: dict[str, int] = {
        "rect": 1,
        "roundRect": 5,
        "ellipse": 9,
        "diamond": 4,
        "triangle": 7,
        "snip1Rect": 88,
        "snip2SameRect": 89,
        "snipRoundRect": 84,
    }

    for shape_def in template.decorative_shapes:
        try:
            shape_type = shape_def.shape_type
            preset = shape_def.shape_preset

            if shape_type == "auto_shape" and preset and preset in _PRESET_ID:
                shape_id = _PRESET_ID[preset]
                shape = slide.shapes.add_shape(
                    shape_id, shape_def.left, shape_def.top, shape_def.width, shape_def.height
                )

                if shape_def.corner_radius is not None:
                    try:
                        shape.adjustments[0] = shape_def.corner_radius
                    except Exception:
                        pass

                if shape_def.fill_color:
                    try:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = _hex_to_rgb(shape_def.fill_color)
                    except Exception:
                        pass

                if shape_def.border_color:
                    try:
                        shape.line.color.rgb = _hex_to_rgb(shape_def.border_color)
                    except Exception:
                        pass

                if shape_def.border_pt:
                    try:
                        shape.line.width = int(shape_def.border_pt * 12700)
                    except Exception:
                        pass

                placed.append(preset)
                log.debug("Placed decorative shape type=%s preset=%s", shape_type, preset)

            elif shape_type == "line":
                connector = slide.shapes.add_connector(1, shape_def.left, shape_def.top,
                                                       shape_def.left + shape_def.width,
                                                       shape_def.top + shape_def.height)
                if shape_def.border_color:
                    try:
                        connector.line.color.rgb = _hex_to_rgb(shape_def.border_color)
                    except Exception:
                        pass

                if shape_def.border_pt:
                    try:
                        connector.line.width = int(shape_def.border_pt * 12700)
                    except Exception:
                        pass

                if shape_def.has_arrowhead:
                    try:
                        connector.line.end_arrow_type = 1  # standard arrow
                    except Exception:
                        pass

                placed.append("line")
                log.debug("Placed line shape with arrowhead=%s", shape_def.has_arrowhead)

        except Exception as e:
            log.warning("Could not place decorative shape %s: %s", shape_def.shape_name, e)

    return placed


# ── Table placement ───────────────────────────────────────────────────────────

def _place_tables(slide, template: SlideTemplate) -> list[str]:
    """Add tables to slide from template table areas."""
    if not template.table_areas:
        return []

    placed = []
    for table_def in template.table_areas:
        try:
            rows = max(table_def.rows, 1)
            cols = max(table_def.columns, 1)

            table_shape = slide.shapes.add_table(rows, cols,
                                                 table_def.left, table_def.top,
                                                 table_def.width, table_def.height)
            table = table_shape.table

            if table_def.header_row:
                for col_idx, header_text in enumerate(table_def.header_row):
                    if col_idx < cols:
                        try:
                            cell = table.cell(0, col_idx)
                            cell.text = header_text
                        except Exception:
                            pass

            placed.append(f"table_{rows}x{cols}")
            log.debug("Placed table %dx%d with %d headers", rows, cols, len(table_def.header_row))

        except Exception as e:
            log.warning("Could not place table: %s", e)

    return placed


# ── execute ────────────────────────────────────────────────────────────────────

def execute(
    presentation_id: str,
    layout: str,
    title: str = "",
    body: str = "",
    speaker_notes: str = "",
    image_overrides: list | None = None,
) -> dict:
    log.debug("add_slide: prs=%s layout=%s", presentation_id, layout)
    prs, guide, _ = _store.get(presentation_id)

    # Load the YAML template for this layout key
    templates = load_slide_templates(guide)
    template = templates.get(layout)
    if template is None:
        available = list(templates.keys())
        log.warning("Unknown layout '%s'. Available: %s", layout, available)
        return {
            "error": f"Layout '{layout}' not found.",
            "available_layouts": available,
            "hint": "Call list_layouts to see all valid keys.",
        }

    # Add slide with the correct PowerPoint layout (for background graphics)
    pptx_layout = _find_pptx_layout(prs, template.layout_name)
    slide = prs.slides.add_slide(pptx_layout)
    log.debug("Added slide with PPTX layout: '%s'", pptx_layout.name)

    # Place decorative shapes (arrows, rectangles, circles, etc.)
    placed_shapes = _place_decorative_shapes(slide, template)

    # Place tables
    placed_tables = _place_tables(slide, template)

    # Place images from template catalog
    placed_images = _place_images(slide, template, guide, image_overrides or [])

    # Fill text areas using exact geometry from the template YAML
    distribution = _distribute_text_by_role(template, title, body)
    filled = []
    for idx, area in enumerate(template.text_areas):
        content = distribution.get(idx, "")
        if not content:
            log.debug("Skipping area[%d] role='%s' (no content)", idx, area.role)
            continue
        _add_text_box(slide, area, content)
        log.debug("Filled area[%d] role='%s' text='%s'", idx, area.role, content[:60])
        filled.append(area.role)

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes

    slide_index = len(prs.slides) - 1
    log.info(
        "Slide %d added: layout='%s' filled=%s shapes=%s tables=%s images=%s prs=%s",
        slide_index, layout, filled, placed_shapes, placed_tables, placed_images, presentation_id,
    )
    return {
        "slide_index": slide_index,
        "layout": layout,
        "filled_areas": filled,
        "placed_shapes": placed_shapes,
        "placed_tables": placed_tables,
        "placed_images": placed_images,
    }
