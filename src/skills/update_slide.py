"""Skill: replace a slide at a given index with new content."""

import logging

from pptx.oxml.ns import qn

from src.skills import _store
from src.skills.add_slide import (
    _add_text_box,
    _distribute_text_by_role,
    _find_pptx_layout,
    _place_decorative_shapes,
    _place_images,
    _place_tables,
)
from src.style_guides.loader import load_slide_templates

log = logging.getLogger(__name__)

TOOL_SPEC = {
    "type": "function",
    "function": {
        "name": "update_slide",
        "description": (
            "Replace a slide at a specific index with improved content. "
            "Use this after render_slides reveals issues such as text overflow, empty areas, "
            "or a poor layout choice. You may change the layout key to a different one. "
            "The slide is rebuilt in place — all other slides are unaffected."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "presentation_id": {"type": "string"},
                "slide_index": {
                    "type": "integer",
                    "description": "0-based index of the slide to replace.",
                },
                "layout": {
                    "type": "string",
                    "description": "Layout key (from list_layouts). Can be different from the original.",
                },
                "title": {"type": "string", "description": "Text for the title area."},
                "body": {
                    "type": "string",
                    "description": (
                        "Text for the body/subtitle area. "
                        "For two-column layouts, separate left and right with \\n---\\n."
                    ),
                },
                "speaker_notes": {"type": "string"},
                "image_overrides": {
                    "type": "array",
                    "description": "Optional image overrides — same format as add_slide.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "role": {"type": "string"},
                            "image_key": {"type": "string"},
                        },
                        "required": ["role", "image_key"],
                    },
                },
            },
            "required": ["presentation_id", "slide_index", "layout"],
        },
    },
}


def _delete_slide(prs, index: int) -> None:
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.get(qn("r:id"))
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)


def _move_slide_to(prs, from_index: int, to_index: int) -> None:
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[from_index]
    sldIdLst.remove(sldId)
    sldIdLst.insert(to_index, sldId)


def execute(
    presentation_id: str,
    slide_index: int,
    layout: str,
    title: str = "",
    body: str = "",
    speaker_notes: str = "",
    image_overrides: list | None = None,
) -> dict:
    log.debug("update_slide: prs=%s index=%d layout=%s", presentation_id, slide_index, layout)
    prs, guide, _ = _store.get(presentation_id)

    slide_count = len(prs.slides)
    if slide_index < 0 or slide_index >= slide_count:
        return {
            "error": f"slide_index {slide_index} out of range (presentation has {slide_count} slides).",
        }

    templates = load_slide_templates(guide)
    template = templates.get(layout)
    if template is None:
        return {
            "error": f"Layout '{layout}' not found.",
            "available_layouts": list(templates.keys()),
            "hint": "Call list_layouts to see all valid keys.",
        }

    # Delete the old slide and append a fresh one, then move it back
    _delete_slide(prs, slide_index)

    pptx_layout = _find_pptx_layout(prs, template.layout_name)
    slide = prs.slides.add_slide(pptx_layout)

    placed_shapes = _place_decorative_shapes(slide, template)
    placed_tables = _place_tables(slide, template)
    placed_images = _place_images(slide, template, guide, image_overrides or [])

    distribution = _distribute_text_by_role(template, title, body)
    filled = []
    for idx, area in enumerate(template.text_areas):
        content = distribution.get(idx, "")
        if not content:
            continue
        _add_text_box(slide, area, content)
        filled.append(area.role)

    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes

    # Move the newly appended slide back to its original position
    new_last = len(prs.slides) - 1
    if new_last != slide_index:
        _move_slide_to(prs, new_last, slide_index)

    log.info(
        "Slide %d updated: layout='%s' filled=%s shapes=%s tables=%s images=%s prs=%s",
        slide_index, layout, filled, placed_shapes, placed_tables, placed_images, presentation_id,
    )
    return {
        "slide_index": slide_index,
        "layout": layout,
        "filled_areas": filled,
        "placed_shapes": placed_shapes,
        "placed_tables": placed_tables,
        "placed_images": placed_images,
    }
