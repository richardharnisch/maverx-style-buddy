"""Tool registry — wraps build-script calls and file I/O as OpenAI function-calling tools."""

import json
import subprocess
from pathlib import Path

_REPO_ROOT = Path(__file__).parents[2]
_SKILL_ROOT = _REPO_ROOT / ".agents" / "skills"
_BUILD_SCRIPT = _SKILL_ROOT / "maverx-presentation-builder" / "scripts" / "build_maverx_artifacts.py"
_TEMPLATE_PPTX = _SKILL_ROOT / "maverx-presentation-builder" / "assets" / "slides" / "templates.pptx"
_OUTPUTS_DIR = _REPO_ROOT / "outputs"
_REQUIRED_LESSON_PLAN_KEYS = [
    "schema_version", "training", "intake_summary", "research_evidence",
    "programme_learning_outcomes", "sessions", "validation",
]
_TEMPLATE_IDS = [
    "01-deck-title", "02-process-slide", "03-unstructured-three-section-slide",
    "04-agenda", "05-text-slide", "06-dark-text-slide",
    "07-complex-layout-slide-1", "08-complex-layout-slide-2", "09-complex-layout-slide-3",
    "10-longer-text-slide", "11-tabular-slide", "12-itemized-text-boxes",
    "13-four-section-slide", "14-dark-section-title-slide", "15-extra-process-timetable-slide",
    "16-three-section-slide", "17-theory-topic-slides", "18-section-title",
    "19-timeline-process-slide", "20-hand-out-slide", "21-big-question",
    "22-break-time", "23-debrief",
]


# ── tool implementations ───────────────────────────────────────────────────────

def _build_presentation(lesson_plan_path: str, out_dir: str | None = None) -> dict:
    cmd = [
        "uv", "run", "python", str(_BUILD_SCRIPT),
        lesson_plan_path,
        "--template", str(_TEMPLATE_PPTX),
        "--clean",
    ]
    if out_dir:
        cmd += ["--out-dir", out_dir]
    result = subprocess.run(cmd, capture_output=True, text=True, cwd=str(_REPO_ROOT))
    if result.returncode != 0:
        return {"success": False, "error": result.stderr.strip(), "stdout": result.stdout.strip()}
    return {"success": True, "manifest": result.stdout.strip()}


_BRAND_FONTS = {"Raleway", "Space Grotesk"}
_BRAND_COLORS = {"0D006A", "F2F2F2", "FFFFFF", "262626"}


def _validate_pptx(pptx_path: str) -> dict:
    from pptx import Presentation  # noqa: PLC0415

    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        return {"valid": False, "error": str(exc)}

    slides_out = []
    for i, slide in enumerate(prs.slides):
        warnings: list[str] = []

        # Speaker notes
        notes_text = ""
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        if not notes_text:
            warnings.append("missing speaker notes")

        shapes_out = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            text = shape.text_frame.text.strip()

            # Resolve font/size/color from the first non-empty run
            font_name = font_size_pt = color_hex = None
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not font_name and run.font.name:
                        font_name = run.font.name
                    if not font_size_pt and run.font.size:
                        font_size_pt = round(run.font.size / 12700)
                    if color_hex is None:
                        try:
                            color_hex = str(run.font.color.rgb)
                        except Exception:
                            pass
                if font_name:
                    break

            # Overflow estimate: chars / (shape width in pt / glyph width) > lines that fit
            overflow_risk = False
            if text and font_size_pt and shape.width and shape.height:
                w_pt = shape.width / 12700
                h_pt = shape.height / 12700
                cpl = max(1, int(w_pt / (font_size_pt * 0.55)))
                lines_fit = max(1, int(h_pt / (font_size_pt * 1.35)))
                est_lines = sum(
                    max(1, len(ln) // cpl + 1) for ln in text.splitlines()
                )
                overflow_risk = est_lines > lines_fit

            if not text:
                warnings.append(f"shape '{shape.name}' is empty")
            if overflow_risk:
                warnings.append(
                    f"shape '{shape.name}' may overflow "
                    f"({len(text)} chars, ~{lines_fit} lines available)"
                )
            if font_name and font_name not in _BRAND_FONTS:
                warnings.append(f"shape '{shape.name}' uses off-brand font '{font_name}'")
            if color_hex and color_hex.upper() not in _BRAND_COLORS:
                warnings.append(f"shape '{shape.name}' uses off-brand color #{color_hex}")

            shapes_out.append({
                "name": shape.name,
                "text": text[:300],
                "char_count": len(text),
                "empty": not text,
                "font": font_name,
                "font_size_pt": font_size_pt,
                "color": f"#{color_hex}" if color_hex else None,
                "overflow_risk": overflow_risk,
            })

        slides_out.append({
            "index": i + 1,
            "layout": slide.slide_layout.name,
            "speaker_notes_preview": notes_text[:300],
            "shapes": shapes_out,
            "warnings": warnings,
        })

    total_warnings = [w for s in slides_out for w in s["warnings"]]
    return {
        "valid": True,
        "pptx_path": pptx_path,
        "slide_count": len(prs.slides),
        "slides": slides_out,
        "summary": {
            "total_warnings": len(total_warnings),
            "slides_with_warnings": sum(1 for s in slides_out if s["warnings"]),
            "empty_shapes": sum(1 for s in slides_out for sh in s["shapes"] if sh["empty"]),
            "overflow_risks": sum(1 for s in slides_out for sh in s["shapes"] if sh["overflow_risk"]),
            "off_brand_issues": sum(
                1 for w in total_warnings if "off-brand" in w
            ),
        },
    }


def _list_templates() -> dict:
    return {"templates": _TEMPLATE_IDS}


def _write_lesson_plan(slug: str, content: str) -> dict:
    try:
        parsed = json.loads(content)
    except json.JSONDecodeError as exc:
        return {"success": False, "error": f"Invalid JSON: {exc}"}
    out_path = _OUTPUTS_DIR / slug / "lesson_plan.json"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(json.dumps(parsed, indent=2, ensure_ascii=False), encoding="utf-8")
    return {"success": True, "path": str(out_path)}


def _validate_lesson_plan(content: str) -> dict:
    try:
        data = json.loads(content)
    except json.JSONDecodeError as exc:
        return {"valid": False, "errors": [f"Invalid JSON: {exc}"]}
    errors = [f"Missing required field: {k}" for k in _REQUIRED_LESSON_PLAN_KEYS if k not in data]
    for session in data.get("sessions", []):
        for key in ["session_n", "title", "deck_outline", "pre_bite", "post_bite"]:
            if key not in session:
                errors.append(f"Session {session.get('session_n', '?')} missing: {key}")
    return {"valid": not errors, "errors": errors} if errors else {"valid": True}


_DISPATCHERS = {
    "build_presentation": _build_presentation,
    "validate_pptx": _validate_pptx,
    "list_templates": _list_templates,
    "write_lesson_plan": _write_lesson_plan,
    "validate_lesson_plan": _validate_lesson_plan,
}

# ── tool specs (OpenAI function-calling format) ────────────────────────────────

_PRESENTATION_TOOLS: list[dict] = [
    {
        "type": "function",
        "function": {
            "name": "build_presentation",
            "description": (
                "Build all PPTX decks and DOCX artifacts from a lesson_plan.json "
                "using the Maverx branded template. Returns the path to manifest.json."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "lesson_plan_path": {
                        "type": "string",
                        "description": "Absolute or relative path to lesson_plan.json",
                    },
                    "out_dir": {
                        "type": "string",
                        "description": "Output directory (defaults to lesson_plan parent / presentation_artifacts)",
                    },
                },
                "required": ["lesson_plan_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "validate_pptx",
            "description": (
                "Inspect a generated PPTX for content and formatting issues. "
                "Returns per-slide shape text, font/color compliance against Maverx brand, "
                "overflow risk estimates, empty shape detection, and speaker note coverage. "
                "Use this after build_presentation to review quality before reporting done."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "pptx_path": {"type": "string", "description": "Path to the .pptx file"},
                },
                "required": ["pptx_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_templates",
            "description": (
                "List all 23 available Maverx slide template IDs. "
                "Consult slide_dictionary.md in the skill assets for guidance on which to choose."
            ),
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
]

_TRAINING_TOOLS: list[dict] = [
    {
        "type": "function",
        "function": {
            "name": "write_lesson_plan",
            "description": (
                "Write the completed lesson_plan.json to disk at "
                "outputs/<slug>/lesson_plan.json. "
                "Call only after the JSON passes validate_lesson_plan."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "slug": {
                        "type": "string",
                        "description": "Kebab-case training slug (e.g. 'lean-six-sigma-intro')",
                    },
                    "content": {
                        "type": "string",
                        "description": "Full lesson_plan.json serialised as a JSON string",
                    },
                },
                "required": ["slug", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "validate_lesson_plan",
            "description": (
                "Check a lesson_plan.json string for required top-level and session-level fields. "
                "Returns {valid: true} or {valid: false, errors: [...]}."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {
                        "type": "string",
                        "description": "The lesson_plan.json content as a JSON string",
                    },
                },
                "required": ["content"],
            },
        },
    },
]


# ── registry ───────────────────────────────────────────────────────────────────

class SkillRegistry:
    def __init__(self, skill: str = "combined") -> None:
        self.skill = skill

    def tool_specs(self) -> list[dict]:
        if self.skill == "training-builder":
            return _TRAINING_TOOLS
        if self.skill == "presentation-builder":
            return _PRESENTATION_TOOLS
        return _TRAINING_TOOLS + _PRESENTATION_TOOLS

    def dispatch(self, name: str, args: dict) -> dict:
        fn = _DISPATCHERS.get(name)
        if fn is None:
            return {"error": f"Unknown tool: {name}"}
        return fn(**args)
