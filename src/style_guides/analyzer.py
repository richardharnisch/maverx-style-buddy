"""Analyze a branded PPTX template and emit one YAML + README per unique slide type.

Run once per template; results are cached in style_guides/{stem}/.
"""

# Set to False to skip LLM calls and use heuristic descriptions instead.
USE_LLM_DESCRIPTIONS = False

# Set to True to send each extracted image to a vision model for a short label.
# Runs once at analysis time; labels are cached in images/catalog.yaml.
# Independent of USE_LLM_DESCRIPTIONS — can be on while descriptions are off.
USE_LLM_IMAGE_LABELS = True

import hashlib
import logging
import os
import re
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.style_guides.schema import DecorativeShape, ImageArea, SlideTemplate, TextArea

log = logging.getLogger(__name__)

# ── helpers ────────────────────────────────────────────────────────────────────

def _slugify(s: str) -> str:
    s = s.lower().strip()
    s = re.sub(r"[^\w\s-]", "", s)
    s = re.sub(r"[\s_-]+", "_", s)
    return s.strip("_") or "slide"


def _fingerprint(slide, W: int, H: int) -> str:
    """Stable hash for a slide's text-box layout.

    Buckets each text-box's (left, top, width) to the nearest 5% of slide
    dimensions, then sorts so order-independent differences don't produce
    different hashes.
    """
    bw = max(W // 20, 1)
    bh = max(H // 20, 1)
    boxes = sorted(
        (
            round(s.left / bw),
            round(s.top / bh),
            round(s.width / bw),
        )
        for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    )
    return f"{slide.slide_layout.name}|{boxes}"


def _extract_color(run) -> str | None:
    """Return '#RRGGBB' if the run has an explicit RGB color, else None."""
    try:
        color_type = run.font.color.type
        if color_type is not None:
            rgb = run.font.color.rgb
            return f"#{str(rgb).upper()}"
    except Exception:
        pass
    return None


def _alignment_name(para) -> str | None:
    """Return alignment string or None if inherited."""
    try:
        al = para.alignment
        if al is not None:
            name = str(al).split(".")[-1].lower()  # PP_ALIGN.CENTER → "center"
            if name in ("left", "center", "right", "justify", "distribute"):
                return name
    except Exception:
        pass
    return None


def _extract_auto_shape_style(shape) -> dict:
    """Extract preset geometry, fill color and corner radius from an auto shape."""
    from pptx.oxml.ns import qn as _qn
    preset = bg_color = corner_radius = None
    try:
        spPr = shape._element.find(_qn("p:spPr"))
        if spPr is not None:
            pg = spPr.find(_qn("a:prstGeom"))
            if pg is not None:
                preset = pg.get("prst")
    except Exception:
        pass
    try:
        fill = shape.fill
        if fill.type is not None:
            bg_color = f"#{fill.fore_color.rgb}"
    except Exception:
        pass
    try:
        adjs = list(shape.adjustments)
        if adjs:
            corner_radius = round(float(adjs[0]), 3)
    except Exception:
        pass
    return {"shape_preset": preset, "bg_color": bg_color, "corner_radius": corner_radius}


def _extract_text_areas(slide) -> list[dict]:
    """Return raw dicts for TEXT_BOX and text-bearing AUTO_SHAPE on the slide."""
    areas = []
    for shape in slide.shapes:
        is_textbox = shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        is_auto = shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        if not (is_textbox or is_auto):
            continue
        if is_auto and (not shape.has_text_frame or not shape.text_frame.text.strip()):
            continue

        font_name = font_size_pt = bold = italic = underline = color = alignment = None
        text_preview = ""
        try:
            if shape.has_text_frame:
                text_preview = shape.text_frame.text[:120]
                for para in shape.text_frame.paragraphs:
                    if alignment is None:
                        alignment = _alignment_name(para)
                    for run in para.runs:
                        if run.font.name and font_name is None:
                            font_name = run.font.name
                        if run.font.size and font_size_pt is None:
                            font_size_pt = round(run.font.size / 12700, 1)
                        if run.font.bold is not None and bold is None:
                            bold = run.font.bold
                        if run.font.italic is not None and italic is None:
                            italic = run.font.italic
                        if run.font.underline is not None and underline is None:
                            underline = run.font.underline
                        if color is None:
                            color = _extract_color(run)
                        break
                    if font_name or font_size_pt:
                        break
        except Exception:
            pass

        style = _extract_auto_shape_style(shape) if is_auto else {"shape_preset": None, "bg_color": None, "corner_radius": None}

        areas.append({
            "shape_name": shape.name,
            "left": int(shape.left or 0),
            "top": int(shape.top or 0),
            "width": int(shape.width or 0),
            "height": int(shape.height or 0),
            "area": int((shape.width or 0) * (shape.height or 0)),
            "font_name": font_name,
            "font_size_pt": font_size_pt,
            "bold": bold,
            "italic": italic,
            "underline": underline,
            "color": color,
            "alignment": alignment,
            "shape_preset": style["shape_preset"],
            "bg_color": style["bg_color"],
            "corner_radius": style["corner_radius"],
            "text_preview": text_preview,
        })
    return areas


def _extract_decorative_shapes(slide, W: int, H: int, captured_names: set) -> tuple[list[dict], int]:
    """Extract non-text, non-image decorative shapes. Returns (shapes, table_count)."""
    from pptx.oxml.ns import qn as _qn
    _DECO_TYPES = {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM}
    min_area = W * H * 0.003  # ignore shapes smaller than 0.3% of slide

    shapes = []
    table_count = 0
    for shape in slide.shapes:
        st = shape.shape_type
        if int(st) == 19:  # TABLE
            table_count += 1
            continue
        if st not in _DECO_TYPES:
            continue
        if shape.name in captured_names:
            continue  # already recorded as a styled text area

        # Skip if has non-empty text (captured as text area)
        try:
            if shape.has_text_frame and shape.text_frame.text.strip():
                continue
        except Exception:
            pass

        left = int(shape.left or 0)
        top = int(shape.top or 0)
        width = int(shape.width or 0)
        height = int(shape.height or 0)

        # Filter out invisible/hairline shapes
        if width * height < min_area and st != MSO_SHAPE_TYPE.LINE:
            continue

        # Preset geometry
        preset = None
        try:
            spPr = shape._element.find(_qn("p:spPr"))
            if spPr is not None:
                pg = spPr.find(_qn("a:prstGeom"))
                if pg is not None:
                    preset = pg.get("prst")
        except Exception:
            pass

        # Fill color
        fill_color = None
        try:
            if shape.fill.type is not None:
                fill_color = f"#{shape.fill.fore_color.rgb}"
        except Exception:
            pass

        # Border
        border_color = border_pt = None
        try:
            ln = shape.line
            if ln.color.type is not None:
                border_color = f"#{ln.color.rgb}"
            if ln.width:
                border_pt = round(ln.width / 12700, 1)
        except Exception:
            pass

        # Skip shapes with no fill and no border (invisible guides)
        if fill_color is None and border_color is None:
            continue

        # Corner radius
        corner_radius = None
        try:
            adjs = list(shape.adjustments)
            if adjs:
                corner_radius = round(float(adjs[0]), 3)
        except Exception:
            pass

        type_name = "line" if st == MSO_SHAPE_TYPE.LINE else ("freeform" if st == MSO_SHAPE_TYPE.FREEFORM else "auto_shape")
        shapes.append({
            "shape_name": shape.name,
            "shape_type": type_name,
            "shape_preset": preset,
            "left": left, "top": top, "width": width, "height": height,
            "fill_color": fill_color,
            "border_color": border_color,
            "border_pt": border_pt,
            "corner_radius": corner_radius,
        })
    return shapes, table_count


def _extract_image_areas(slide, W: int, H: int) -> list[dict]:
    """Return raw dicts for every PICTURE shape on the slide."""
    areas = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        left = int(shape.left or 0)
        top = int(shape.top or 0)
        width = int(shape.width or 0)
        height = int(shape.height or 0)

        # Assign role by size relative to slide
        slide_area = W * H
        shape_area = width * height
        coverage = shape_area / slide_area if slide_area else 0
        if coverage > 0.35:
            role = "background"
        elif coverage < 0.04:
            role = "icon"
        else:
            role = "photo"

        try:
            blob = shape.image.blob
            content_type = shape.image.content_type
            ext = shape.image.ext
        except Exception:
            blob = content_type = ext = None

        areas.append({
            "shape_name": shape.name,
            "role": role,
            "left": left,
            "top": top,
            "width": width,
            "height": height,
            "blob": blob,
            "content_type": content_type,
            "ext": ext or "png",
        })
    return areas


def _save_images(raw_image_areas: list[dict], images_dir: Path, catalog: dict) -> list[dict]:
    """Dedup images by content hash, save to images_dir, update catalog.

    Returns image area dicts with image_key set (blob removed).
    """
    images_dir.mkdir(parents=True, exist_ok=True)
    result = []
    for raw in raw_image_areas:
        blob = raw.get("blob")
        image_key = None
        if blob:
            h = hashlib.sha256(blob).hexdigest()[:12]
            ext = raw.get("ext", "png")
            image_key = f"img_{h}"
            img_path = images_dir / f"{image_key}.{ext}"
            if not img_path.exists():
                img_path.write_bytes(blob)
            entry = catalog.setdefault(image_key, {
                "key": image_key,
                "filename": img_path.name,
                "content_type": raw.get("content_type", "image/png"),
                "source_slides": [],
                "usage_count": 0,
                "description": None,
            })
            entry["usage_count"] += 1
        result.append({
            "shape_name": raw["shape_name"],
            "role": raw["role"],
            "left": raw["left"],
            "top": raw["top"],
            "width": raw["width"],
            "height": raw["height"],
            "image_key": image_key,
        })
    return result


def _assign_roles(areas: list[dict], W: int, H: int) -> list[dict]:
    """Add a 'role' key to each area dict in-place (returns same list)."""
    if not areas:
        return areas

    by_top = sorted(areas, key=lambda a: (a["top"], a["left"]))

    # Detect two-column: two boxes at similar vertical position, wide horizontal gap
    col_pair = None
    for i in range(len(by_top) - 1):
        a, b = by_top[i], by_top[i + 1]
        if abs(a["top"] - b["top"]) < H * 0.15 and abs(a["left"] - b["left"]) > W * 0.25:
            col_pair = (i, i + 1)
            break

    if col_pair:
        ci, cj = col_pair
        title_assigned = False
        for k, area in enumerate(by_top):
            if k == ci:
                area["role"] = "left_column"
            elif k == cj:
                area["role"] = "right_column"
            elif not title_assigned and area["top"] < by_top[ci]["top"]:
                area["role"] = "title"
                title_assigned = True
            else:
                area["role"] = "caption"
        return by_top

    # Standard: topmost → title; next → body or subtitle; rest → caption
    for k, area in enumerate(by_top):
        if area["top"] > H * 0.85:
            area["role"] = "caption"
        elif k == 0:
            area["role"] = "title"
        elif k == 1:
            prev_size = by_top[0].get("font_size_pt") or 0
            this_size = area.get("font_size_pt") or 0
            if prev_size and this_size and this_size >= prev_size * 0.6 and area["area"] < by_top[0]["area"] * 2:
                area["role"] = "subtitle"
            else:
                area["role"] = "body"
        else:
            area["role"] = "body"
    return by_top


def _slide_all_text(slide) -> str:
    """Concatenate all text from a slide's shapes for LLM input."""
    parts = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
        except Exception:
            pass
    return "\n".join(parts)[:600]


def _shape_summary(areas: list[dict], image_areas: list[dict], deco_shapes: list[dict], table_count: int, non_text: int) -> str:
    role_list = ", ".join(a["role"] for a in areas) if areas else "none"
    sizes = [f"{a['font_size_pt']}pt" for a in areas if a.get("font_size_pt")]
    size_str = f" (font sizes: {', '.join(sizes)})" if sizes else ""
    img_str = f"; {len(image_areas)} image(s): {', '.join(a['role'] for a in image_areas)}" if image_areas else ""
    deco_presets = [d["shape_preset"] or d["shape_type"] for d in deco_shapes]
    deco_str = f"; decorative shapes: {', '.join(deco_presets)}" if deco_shapes else ""
    table_str = f"; {table_count} table(s)" if table_count else ""
    return f"{len(areas)} text area(s) with roles [{role_list}]{size_str}{img_str}{deco_str}{table_str}; {non_text} other shape(s)"


def _is_light_color(hex_color: str) -> bool:
    """Return True if the hex color is light (luminance > 0.6)."""
    try:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
        return luminance > 0.6
    except Exception:
        return False


def _heuristic_description(areas: list[dict], image_areas: list[dict], deco_shapes: list[dict], table_count: int, non_text: int, layout_name: str) -> str:
    roles = {a["role"] for a in areas}
    img_roles = {a["role"] for a in image_areas}

    # Determine visual tone from text colors
    light_text = any(_is_light_color(a["color"]) for a in areas if a.get("color"))
    has_bg_image = "background" in img_roles
    has_photo = "photo" in img_roles
    has_icon = "icon" in img_roles

    # Detect brand accent colors on title text
    title_color = next((a["color"] for a in areas if a.get("role") == "title" and a.get("color")), None)

    # Visual context prefix
    if has_bg_image and light_text:
        visual = "Dark photo-background slide with white/light text."
    elif has_bg_image:
        visual = "Slide with a large background photo."
    elif has_photo and light_text:
        visual = "Slide with a content photo and light-colored text."
    elif has_photo:
        visual = "Slide with a content photo."
    elif light_text and not areas:
        return (
            f"Full-image or decorative slide using the '{layout_name}' layout. "
            "Use for visual impact, section dividers, or image-only pages."
        )
    elif not areas:
        return (
            f"Full-image or decorative slide using the '{layout_name}' layout. "
            "Use for visual impact, section dividers, or image-only pages."
        )
    elif light_text:
        visual = "Dark-background slide with light/white text."
    elif title_color:
        visual = f"Light slide with branded accent color on title text ({title_color})."
    else:
        visual = ""

    icon_note = " Includes decorative icons." if has_icon else ""

    # Decorative shape notes
    deco_notes = []
    if table_count:
        deco_notes.append(f"Contains a {'table' if table_count == 1 else f'{table_count} tables'}.")
    ovals = [d for d in deco_shapes if d.get("shape_preset") == "ellipse"]
    rects = [d for d in deco_shapes if d.get("shape_preset") in ("rect", "roundRect")]
    lines = [d for d in deco_shapes if d.get("shape_type") == "line"]
    if ovals:
        colors = list({d["fill_color"] for d in ovals if d.get("fill_color")})
        color_hint = f" ({', '.join(colors[:3])})" if colors else ""
        deco_notes.append(f"Has {len(ovals)} colored circle/oval element(s){color_hint}.")
    if rects:
        colors = list({d["fill_color"] for d in rects if d.get("fill_color")})
        color_hint = f" ({', '.join(colors[:3])})" if colors else ""
        deco_notes.append(f"Has {len(rects)} accent rectangle/block(s){color_hint}.")
    if lines:
        deco_notes.append(f"Has {len(lines)} decorative line(s).")
    deco_note = " ".join(deco_notes)

    # Text structure
    if "left_column" in roles and "right_column" in roles:
        prefix = "with a title, " if "title" in roles else ""
        structure = f"Two-column layout {prefix}suited for side-by-side comparisons or split content."
    elif "title" in roles and "subtitle" in roles and len(roles) <= 3:
        title_area = next((a for a in areas if a["role"] == "title"), None)
        font_hint = ""
        if title_area and title_area.get("font_size_pt"):
            font_hint = f" ({title_area['font_size_pt']}pt heading)"
        structure = f"Cover/title slide{font_hint} with a large heading and subtitle. Use for opening slides or major section dividers."
    elif "title" in roles and ("body" in roles or len(areas) > 2):
        body_count = sum(1 for a in areas if a["role"] == "body")
        if body_count > 1:
            structure = f"Content slide with title and {body_count} body text areas. Suited for multi-point or structured information."
        else:
            structure = "Standard content slide with a title and body text area. Use for most informational slides."
    elif "title" in roles and len(roles) == 1:
        structure = "Title-only slide. Use for section headers or impactful transition slides."
    else:
        structure = f"Content slide with {len(areas)} text area(s). Layout: '{layout_name}'."

    parts = [p for p in [visual, structure, icon_note.strip(), deco_note] if p]
    return " ".join(parts)


def _llm_label_image(img_path: Path) -> str | None:
    """Send an image to the configured vision model and return a short label. Returns None on failure."""
    # load_dotenv() before checking env so the key is present even when called directly
    from dotenv import load_dotenv
    load_dotenv()
    api_key = os.environ.get("OPENROUTER_API_KEY")
    if not api_key:
        return None
    try:
        import base64
        from src.agent.client import OpenRouterClient

        blob = img_path.read_bytes()
        b64 = base64.b64encode(blob).decode()
        suffix_map = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", ".gif": "image/gif", ".webp": "image/webp"}
        mime = suffix_map.get(img_path.suffix.lower(), "image/png")
        data_url = f"data:{mime};base64,{b64}"

        # Use default model from env (supports vision on this account)
        client = OpenRouterClient()
        response = client.chat([
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "Describe this image in 8 words or fewer. "
                            "State what it is (e.g. 'company logo', 'office team photo', "
                            "'abstract geometric icon', 'dark texture overlay', 'checkmark icon'). "
                            "Output ONLY the description, nothing else."
                        ),
                    },
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }
        ])
        label = response.choices[0].message.content.strip().strip('"').strip("'")
        log.debug("Image label for '%s': %s", img_path.name, label)
        return label
    except Exception as e:
        log.warning("Image labeling failed for '%s': %s", img_path.name, e)
        return None


def _llm_description(layout_name: str, slide_text: str, shape_summary: str) -> str | None:
    """Call the LLM to generate a concise description. Returns None on failure."""
    api_key = os.environ.get("OPENROUTER_API_KEY")
    if not api_key:
        return None
    try:
        from src.agent.client import OpenRouterClient
        client = OpenRouterClient()
        prompt = (
            f"You are analysing a PowerPoint slide template to help an AI assistant pick the right layout.\n\n"
            f"Layout name: {layout_name}\n"
            f"Shape structure: {shape_summary}\n"
            f"Slide text content:\n---\n{slide_text}\n---\n\n"
            f"Write ONE to TWO concise sentences (max 40 words total) that describe:\n"
            f"1. The visual style (dark/light, photo background, colors, two-column, etc.)\n"
            f"2. What type of content this slide is designed for\n"
            f"3. When an AI presentation builder should choose it\n\n"
            f"Output only the sentence(s), no preamble."
        )
        response = client.chat([{"role": "user", "content": prompt}])
        return response.choices[0].message.content.strip()
    except Exception as e:
        log.warning("LLM description failed: %s", e)
        return None


# ── main public function ───────────────────────────────────────────────────────

def analyze(pptx_path: Path, output_dir: Path) -> list[SlideTemplate]:
    """Inspect every slide in pptx_path, group by structure, write YAMLs + README."""
    log.info("Analysing template: %s → %s", pptx_path.name, output_dir)
    prs = Presentation(str(pptx_path))
    W = int(prs.slide_width)
    H = int(prs.slide_height)
    output_dir.mkdir(parents=True, exist_ok=True)
    images_dir = output_dir / "images"
    image_catalog: dict = {}

    # ── group slides by structural fingerprint ─────────────────────────────────
    groups: dict[str, dict] = {}   # fingerprint → {slide_index, slide, layout_name}
    for idx, slide in enumerate(prs.slides):
        fp = _fingerprint(slide, W, H)
        if fp not in groups:
            groups[fp] = {"slide_index": idx, "slide": slide, "layout_name": slide.slide_layout.name}

    log.info("Found %d unique slide structure(s) across %d slides", len(groups), len(prs.slides))

    # ── assign keys (slug + dedup suffix) ─────────────────────────────────────
    slug_counts: dict[str, int] = {}
    ordered: list[dict] = []
    for fp, g in groups.items():
        slug = _slugify(g["layout_name"])
        slug_counts[slug] = slug_counts.get(slug, 0) + 1
        g["slug"] = slug
        g["fingerprint"] = fp
        ordered.append(g)

    slug_seen: dict[str, int] = {}
    for g in ordered:
        slug = g["slug"]
        slug_seen[slug] = slug_seen.get(slug, 0) + 1
        if slug_counts[slug] == 1:
            g["key"] = slug
        else:
            g["key"] = f"{slug}_{slug_seen[slug]}"

    # ── build SlideTemplate objects and write YAMLs ────────────────────────────
    templates: list[SlideTemplate] = []
    for g in ordered:
        slide = g["slide"]
        layout_name = g["layout_name"]
        key = g["key"]

        raw_areas = _extract_text_areas(slide)
        raw_image_areas = _extract_image_areas(slide, W, H)
        captured_names = {a["shape_name"] for a in raw_areas}
        raw_deco, table_count = _extract_decorative_shapes(slide, W, H, captured_names)
        non_text = sum(
            1 for s in slide.shapes
            if s.shape_type not in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PICTURE)
            and s.name not in captured_names
            and s.name not in {d["shape_name"] for d in raw_deco}
        )
        areas = _assign_roles(raw_areas, W, H)
        image_areas = _save_images(raw_image_areas, images_dir, image_catalog)
        slide_text = _slide_all_text(slide)
        summary = _shape_summary(areas, image_areas, raw_deco, table_count, non_text)

        log.debug("Layout '%s' (key=%s): %s", layout_name, key, summary)

        if USE_LLM_DESCRIPTIONS:
            desc = _llm_description(layout_name, slide_text, summary)
        else:
            desc = None
        if desc:
            log.debug("LLM description for '%s': %s", key, desc)
        else:
            desc = _heuristic_description(areas, image_areas, raw_deco, table_count, non_text, layout_name)
            log.debug("Heuristic description for '%s': %s", key, desc)

        text_area_models = [
            TextArea(
                role=a["role"],
                shape_name=a["shape_name"],
                left=a["left"],
                top=a["top"],
                width=a["width"],
                height=a["height"],
                font_name=a.get("font_name"),
                font_size_pt=a.get("font_size_pt"),
                bold=a.get("bold"),
                italic=a.get("italic"),
                underline=a.get("underline"),
                color=a.get("color"),
                alignment=a.get("alignment"),
                shape_preset=a.get("shape_preset"),
                bg_color=a.get("bg_color"),
                corner_radius=a.get("corner_radius"),
            )
            for a in areas
        ]

        image_area_models = [
            ImageArea(
                role=ia["role"],
                shape_name=ia["shape_name"],
                left=ia["left"],
                top=ia["top"],
                width=ia["width"],
                height=ia["height"],
                image_key=ia.get("image_key"),
            )
            for ia in image_areas
        ]

        deco_models = [
            DecorativeShape(
                shape_name=d["shape_name"],
                shape_type=d["shape_type"],
                shape_preset=d.get("shape_preset"),
                left=d["left"], top=d["top"], width=d["width"], height=d["height"],
                fill_color=d.get("fill_color"),
                border_color=d.get("border_color"),
                border_pt=d.get("border_pt"),
                corner_radius=d.get("corner_radius"),
            )
            for d in raw_deco
        ]

        template = SlideTemplate(
            key=key,
            source_slide_index=g["slide_index"],
            layout_name=layout_name,
            description=desc,
            text_areas=text_area_models,
            image_areas=image_area_models,
            decorative_shapes=deco_models,
            table_count=table_count,
            non_text_shapes=non_text,
        )
        templates.append(template)

        # Write YAML
        yaml_path = output_dir / f"{key}.yaml"
        data = template.model_dump()
        with yaml_path.open("w") as f:
            yaml.dump(data, f, default_flow_style=False, allow_unicode=True, sort_keys=False)
        log.debug("Wrote %s", yaml_path.name)

    # ── label images with vision LLM ──────────────────────────────────────────
    if image_catalog and USE_LLM_IMAGE_LABELS:
        log.info("Labelling %d image(s) with vision model...", len(image_catalog))
        for entry in image_catalog.values():
            if entry.get("description"):
                continue  # already labelled (won't happen on fresh analysis, but useful on incremental runs)
            img_path = images_dir / entry["filename"]
            label = _llm_label_image(img_path)
            entry["description"] = label or "unlabelled"

    # ── write image catalog ────────────────────────────────────────────────────
    if image_catalog:
        catalog_path = images_dir / "catalog.yaml"
        catalog_list = sorted(image_catalog.values(), key=lambda e: -e["usage_count"])
        with catalog_path.open("w") as f:
            yaml.dump(catalog_list, f, default_flow_style=False, allow_unicode=True, sort_keys=False)
        log.info("Wrote image catalog: %d unique image(s)", len(image_catalog))

    # ── write README.md ────────────────────────────────────────────────────────
    _write_readme(output_dir, templates, pptx_path.stem)
    log.info("Analysis complete: %d layout(s) written to %s", len(templates), output_dir)
    return templates


def _write_readme(output_dir: Path, templates: list[SlideTemplate], guide_name: str) -> None:
    lines = [
        f"# {guide_name} — Slide Layouts\n",
        "Use the `key` values below as the `layout` parameter in `add_slide`.\n",
        "---\n",
    ]
    for t in templates:
        font_info = []
        for ta in t.text_areas:
            parts = [f"`{ta.role}`"]
            if ta.font_name:
                parts.append(ta.font_name)
            if ta.font_size_pt:
                parts.append(f"{ta.font_size_pt}pt")
            if ta.bold:
                parts.append("bold")
            if ta.color:
                parts.append(f"color={ta.color}")
            if ta.alignment:
                parts.append(f"align={ta.alignment}")
            font_info.append(" — ".join(parts))

        img_info = [f"`{ia.role}` ({ia.image_key or 'no default'})" for ia in t.image_areas]

        lines += [
            f"## {t.key}\n",
            f"**When to use:** {t.description}\n",
            f"**PowerPoint layout:** `{t.layout_name}`  \n",
            f"**Source slide:** {t.source_slide_index}  \n",
        ]
        if font_info:
            lines.append("**Text areas:**\n")
            for fi in font_info:
                lines.append(f"- {fi}\n")
        else:
            lines.append("**Text areas:** none (image/decorative slide)\n")
        if img_info:
            lines.append("**Image areas:**\n")
            for ii in img_info:
                lines.append(f"- {ii}\n")
        if t.decorative_shapes:
            lines.append("**Decorative shapes:**\n")
            for ds in t.decorative_shapes:
                parts = [ds.shape_preset or ds.shape_type]
                if ds.fill_color:
                    parts.append(f"fill={ds.fill_color}")
                if ds.corner_radius:
                    parts.append(f"radius={ds.corner_radius}")
                if ds.border_color:
                    parts.append(f"border={ds.border_color}")
                lines.append(f"- {' '.join(parts)}\n")
        if t.table_count:
            lines.append(f"**Tables:** {t.table_count}\n")
        lines.append("\n---\n")

    readme_path = output_dir / "README.md"
    with readme_path.open("w") as f:
        f.writelines(lines)
    log.info("Wrote %s", readme_path)
