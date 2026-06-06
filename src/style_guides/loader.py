"""Load and validate a style guide from the /style_guides directory.

Resolution order:
  1. {name}.yaml exists  → load it (may have been auto-generated or hand-edited)
  2. {name}.pptx exists  → extract brand data, write {name}.yaml, return guide
  3. Neither             → FileNotFoundError
"""

import logging
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.oxml.ns import qn

import yaml as _yaml_mod  # alias used in load_slide_templates
from src.style_guides.schema import Colors, Fonts, SlideSize, SlideTemplate, StyleGuide, TextArea

log = logging.getLogger(__name__)

PROJECT_ROOT = Path(__file__).parents[2]
STYLE_GUIDES_DIR = PROJECT_ROOT / "style_guides"

_THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"


# ---------------------------------------------------------------------------
# Extraction helpers
# ---------------------------------------------------------------------------

def _srgb(el) -> str | None:
    """Return '#RRGGBB' from an <a:srgbClr> or <a:sysClr> child, or None."""
    if el is None:
        return None
    node = el.find(qn("a:srgbClr"))
    if node is not None:
        return f"#{node.get('val', '000000').upper()}"
    node = el.find(qn("a:sysClr"))
    if node is not None:
        return f"#{node.get('lastClr', '000000').upper()}"
    return None


def _pt(emu_or_none) -> int | None:
    """Convert EMU to whole points, or return None."""
    if emu_or_none is None:
        return None
    return round(emu_or_none / 12700)


def _extract_from_pptx(pptx_path: Path) -> StyleGuide:
    log.info("Extracting style guide from PPTX: %s", pptx_path.name)
    prs = Presentation(str(pptx_path))
    master = prs.slide_master

    # --- Slide size ---
    slide_size = SlideSize(
        width_emu=int(prs.slide_width),
        height_emu=int(prs.slide_height),
    )

    # --- Theme element (fonts + color scheme) ---
    theme_el = None
    try:
        from lxml import etree
        theme_part = master.part.part_related_by(_THEME_REL)
        theme_el = etree.fromstring(theme_part.blob)
        log.debug("Theme element loaded")
    except Exception:
        log.warning("Could not load theme element — using defaults")

    # --- Fonts from theme font scheme ---
    heading_font = "Calibri"
    body_font = "Calibri"
    if theme_el is not None:
        major = theme_el.find(".//" + qn("a:majorFont"))
        minor = theme_el.find(".//" + qn("a:minorFont"))
        if major is not None:
            latin = major.find(qn("a:latin"))
            if latin is not None:
                heading_font = latin.get("typeface", heading_font)
        if minor is not None:
            latin = minor.find(qn("a:latin"))
            if latin is not None:
                body_font = latin.get("typeface", body_font)

    # --- Font sizes from master placeholders ---
    heading_size = 32
    body_size = 18
    for ph in master.placeholders:
        idx = ph.placeholder_format.idx
        try:
            for para in ph.text_frame.paragraphs:
                sz = para.runs[0].font.size if para.runs else None
                if sz and idx == 0:
                    heading_size = _pt(sz) or heading_size
                elif sz and idx == 1:
                    body_size = _pt(sz) or body_size
        except Exception:
            pass

    fonts = Fonts(
        heading=heading_font,
        body=body_font,
        heading_size=heading_size,
        body_size=body_size,
    )
    log.debug("Fonts: heading=%s %dpt, body=%s %dpt", heading_font, heading_size, body_font, body_size)

    # --- Colors from theme color scheme ---
    primary = "#1F3864"
    secondary = "#2E75B6"
    accent = "#ED7D31"
    background = "#FFFFFF"
    text = "#212121"

    if theme_el is not None:
        clr_scheme = theme_el.find(".//" + qn("a:clrScheme"))
        if clr_scheme is not None:
            primary   = _srgb(clr_scheme.find(qn("a:accent1"))) or primary
            secondary = _srgb(clr_scheme.find(qn("a:accent2"))) or secondary
            accent    = _srgb(clr_scheme.find(qn("a:accent3"))) or accent
            text      = _srgb(clr_scheme.find(qn("a:dk1")))     or text
            background = _srgb(clr_scheme.find(qn("a:lt1")))    or background

    # Background from master solid fill takes priority over theme lt1
    try:
        fill = master.background.fill
        if fill.type is not None:
            rgb = fill.fore_color.rgb
            background = f"#{rgb}"
    except Exception:
        pass

    colors = Colors(
        primary=primary,
        secondary=secondary,
        accent=accent,
        background=background,
        text=text,
    )
    log.debug("Colors: primary=%s secondary=%s background=%s text=%s", primary, secondary, background, text)

    return StyleGuide(
        name=pptx_path.stem,
        template_pptx=pptx_path.name,
        colors=colors,
        fonts=fonts,
        slide_size=slide_size,
    )


def _write_yaml(guide: StyleGuide, yaml_path: Path) -> None:
    data = guide.model_dump(exclude_none=False)
    with yaml_path.open("w") as f:
        yaml.dump(data, f, default_flow_style=False, allow_unicode=True, sort_keys=False)
    log.info("Generated %s", yaml_path.name)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def _ensure_layouts(guide: StyleGuide, pptx_path: Path) -> StyleGuide:
    """Run the analyzer if the layouts directory doesn't exist yet."""
    layouts_dir = STYLE_GUIDES_DIR / pptx_path.stem
    if not layouts_dir.exists():
        log.info("No layout cache found — running template analysis for '%s'", pptx_path.name)
        from src.style_guides.analyzer import analyze
        analyze(pptx_path, layouts_dir)
    rel = str(layouts_dir.relative_to(PROJECT_ROOT))
    return guide.model_copy(update={"layouts_dir": rel})


def load_style_guide(name: str) -> StyleGuide:
    log.debug("Loading style guide '%s'", name)
    yaml_path = STYLE_GUIDES_DIR / f"{name}.yaml"
    pptx_path = STYLE_GUIDES_DIR / f"{name}.pptx"

    if yaml_path.exists():
        log.debug("Found %s", yaml_path.name)
        with yaml_path.open() as f:
            data = yaml.safe_load(f)
        guide = StyleGuide(**data)
        if guide.template_pptx is None and pptx_path.exists():
            guide = guide.model_copy(update={"template_pptx": pptx_path.name})
            log.debug("Auto-wired template: %s", pptx_path.name)
    elif pptx_path.exists():
        guide = _extract_from_pptx(pptx_path)
        _write_yaml(guide, yaml_path)
    else:
        raise FileNotFoundError(
            f"No style guide found for '{name}'. "
            f"Drop '{name}.pptx' into the style_guides/ directory."
        )

    # Ensure slide layout templates are analysed and linked
    resolved = resolve_template_path(guide)
    if resolved is not None:
        guide = _ensure_layouts(guide, resolved)

    return guide


def resolve_template_path(guide: StyleGuide) -> Path | None:
    if guide.template_pptx is None:
        return None
    path = STYLE_GUIDES_DIR / guide.template_pptx
    return path if path.exists() else None


def load_slide_templates(guide: StyleGuide) -> dict[str, SlideTemplate]:
    """Load all YAML layout files from guide.layouts_dir, keyed by slug."""
    if not guide.layouts_dir:
        return {}
    layouts_dir = PROJECT_ROOT / guide.layouts_dir
    if not layouts_dir.exists():
        return {}
    templates: dict[str, SlideTemplate] = {}
    for yaml_file in sorted(layouts_dir.glob("*.yaml")):
        try:
            with yaml_file.open() as f:
                data = _yaml_mod.safe_load(f)
            templates[data["key"]] = SlideTemplate(**data)
        except Exception as e:
            log.warning("Could not load layout '%s': %s", yaml_file.name, e)
    return templates
