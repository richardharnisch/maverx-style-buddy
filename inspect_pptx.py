"""Inspect every python-pptx-accessible parameter in a .pptx file.

Usage:
    uv run python inspect_pptx.py <file.pptx> [--slides] [--xml]

Flags:
    --slides    Also dump each content slide (skipped by default for large files)
    --xml       Print raw XML for each shape/placeholder
"""

import sys
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from rich import box
from rich.console import Console
from rich.panel import Panel
from rich.table import Table
from rich.tree import Tree

console = Console()
DUMP_SLIDES = "--slides" in sys.argv
DUMP_XML = "--xml" in sys.argv


# ── Formatting helpers ─────────────────────────────────────────────────────────

def emu(val) -> str:
    if val is None:
        return "None"
    inches = val / 914400
    return f"{val} EMU  ({inches:.3f}\")"

def pt(val) -> str:
    if val is None:
        return "None"
    return f"{val / 12700:.1f} pt"

def safe(fn):
    try:
        v = fn()
        return str(v) if v is not None else "[dim]None[/dim]"
    except Exception as e:
        return f"[red]ERR: {e}[/red]"


# ── Color ──────────────────────────────────────────────────────────────────────

def color_str(color_obj) -> str:
    try:
        t = color_obj.type
        if t is None:
            return "[dim]inherited[/dim]"
        if str(t) in ("RGB (1)", "1"):
            return f"RGB #{color_obj.rgb}"
        if "THEME" in str(t):
            tc = color_obj.theme_color
            bright = color_obj.brightness
            return f"THEME {tc.name}  brightness={bright:.2f}"
        return str(t)
    except Exception as e:
        return f"[red]ERR:{e}[/red]"


# ── Fill ───────────────────────────────────────────────────────────────────────

def fill_tree(fill, parent: Tree) -> None:
    try:
        ft = fill.type
        b = parent.add(f"[cyan]fill.type[/cyan] = {ft}")
        if ft is not None and "SOLID" in str(ft):
            b.add(f"[cyan]fill.fore_color[/cyan] = {color_str(fill.fore_color)}")
    except Exception as e:
        parent.add(f"[red]fill ERR: {e}[/red]")


# ── Font ───────────────────────────────────────────────────────────────────────

def font_tree(font, parent: Tree) -> None:
    fields = [
        ("name", lambda: font.name),
        ("size", lambda: pt(font.size)),
        ("bold", lambda: font.bold),
        ("italic", lambda: font.italic),
        ("underline", lambda: font.underline),
        ("color", lambda: color_str(font.color)),
    ]
    for name, fn in fields:
        parent.add(f"[cyan]font.{name}[/cyan] = {safe(fn)}")


# ── Text frame ────────────────────────────────────────────────────────────────

def tf_tree(tf, parent: Tree, max_paras: int = 3) -> None:
    parent.add(f"[cyan]word_wrap[/cyan] = {safe(lambda: tf.word_wrap)}")
    parent.add(f"[cyan]auto_size[/cyan] = {safe(lambda: tf.auto_size)}")
    parent.add(f"[cyan]text (first 120)[/cyan] = {repr(tf.text[:120])}")
    for i, para in enumerate(tf.paragraphs[:max_paras]):
        pb = parent.add(f"[green]paragraph[{i}][/green]  runs={len(para.runs)}  align={safe(lambda: para.alignment)}")
        if para.runs:
            rb = pb.add(f"[green]run[0][/green]  text={repr(para.runs[0].text[:60])}")
            font_tree(para.runs[0].font, rb)
    if len(tf.paragraphs) > max_paras:
        parent.add(f"[dim]… {len(tf.paragraphs) - max_paras} more paragraph(s)[/dim]")


# ── Shape ─────────────────────────────────────────────────────────────────────

def shape_tree(shape, parent: Tree, show_xml: bool = False) -> None:
    node = parent.add(
        f"[bold yellow]shape[/bold yellow]  id={shape.shape_id}  name={repr(shape.name)}  "
        f"type={safe(lambda: shape.shape_type.name)}"
    )
    node.add(f"left={emu(shape.left)}  top={emu(shape.top)}")
    node.add(f"width={emu(shape.width)}  height={emu(shape.height)}")

    # Placeholder info
    if shape.is_placeholder:
        phf = shape.placeholder_format
        ph = node.add(f"[magenta]placeholder_format[/magenta]  idx={phf.idx}  type={phf.type}")
        ph.add(f"[cyan]name[/cyan] = {repr(shape.name)}")

    # Fill
    try:
        fill_tree(shape.fill, node)
    except Exception:
        pass

    # Line
    try:
        line = shape.line
        ln = node.add("[cyan]line[/cyan]")
        ln.add(f"width={safe(lambda: pt(line.width))}  color={color_str(line.color)}")
    except Exception:
        pass

    # Text frame
    if shape.has_text_frame:
        tfn = node.add("[cyan]text_frame[/cyan]")
        tf_tree(shape.text_frame, tfn)

    # Raw XML
    if show_xml:
        xml = etree.tostring(shape._element, pretty_print=True).decode()
        node.add(f"[dim]{xml[:800]}[/dim]")


# ── Slide layout ──────────────────────────────────────────────────────────────

def layout_tree(layout, idx: int) -> Tree:
    root = Tree(f"[bold blue]layout[{idx}][/bold blue]  name={repr(layout.name)}")

    # Placeholders table
    phs = list(layout.placeholders)
    if phs:
        t = Table("idx", "type", "name", "left", "top", "width", "height", box=box.SIMPLE_HEAVY)
        for ph in phs:
            phf = ph.placeholder_format
            t.add_row(
                str(phf.idx),
                phf.type.name,
                ph.name,
                emu(ph.left),
                emu(ph.top),
                emu(ph.width),
                emu(ph.height),
            )
        root.add(Panel(t, title="placeholders", border_style="blue"))
    else:
        root.add("[dim]no placeholders[/dim]")

    # All shapes (non-placeholder too)
    shapes = list(layout.shapes)
    root.add(f"[cyan]total shapes[/cyan] = {len(shapes)}")
    shapes_node = root.add("[cyan]shapes[/cyan]")
    for s in shapes:
        shape_tree(s, shapes_node, show_xml=DUMP_XML)

    # Background fill
    try:
        bg_node = root.add("[cyan]background.fill[/cyan]")
        fill_tree(layout.background.fill, bg_node)
    except Exception:
        pass

    return root


# ── Slide master ──────────────────────────────────────────────────────────────

def master_tree(master) -> Tree:
    root = Tree("[bold red]slide_master[/bold red]")

    # Theme colors via XML
    try:
        _THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        theme_part = next(
            rel.target_part
            for rel in master.part.rels.values()
            if rel.reltype == _THEME_REL
        )
        xml = etree.fromstring(theme_part.blob)
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

        colors_node = root.add("[bold]theme colors[/bold]")
        scheme = xml.find(f".//{{{ns}}}clrScheme")
        if scheme is not None:
            for child in scheme:
                tag = child.tag.split("}")[-1]
                color_child = child[0] if len(child) else None
                if color_child is not None:
                    val = color_child.get("val") or color_child.get("lastClr") or color_child.get("val", "?")
                    colors_node.add(f"[cyan]{tag}[/cyan] = {color_child.tag.split('}')[-1]}  val={val}")

        fonts_node = root.add("[bold]theme fonts[/bold]")
        fmt = xml.find(f".//{{{ns}}}fontScheme")
        if fmt is not None:
            for family in fmt:
                tag = family.tag.split("}")[-1]
                latin = family.find(f"{{{ns}}}latin")
                ea = family.find(f"{{{ns}}}ea")
                cs = family.find(f"{{{ns}}}cs")
                fnode = fonts_node.add(f"[cyan]{tag}[/cyan]")
                if latin is not None:
                    fnode.add(f"latin = {latin.get('typeface')}")
                if ea is not None:
                    fnode.add(f"ea    = {ea.get('typeface')}")
                if cs is not None:
                    fnode.add(f"cs    = {cs.get('typeface')}")
    except Exception as e:
        root.add(f"[red]theme extraction failed: {e}[/red]")

    # Master placeholders
    phs = list(master.placeholders)
    root.add(f"[cyan]placeholders[/cyan] = {len(phs)}")
    for ph in phs:
        phf = ph.placeholder_format
        root.add(f"  idx={phf.idx}  type={phf.type.name}  name={repr(ph.name)}  w={emu(ph.width)}  h={emu(ph.height)}")

    # Master shapes
    root.add(f"[cyan]shapes[/cyan] = {len(list(master.shapes))}")
    shapes_node = root.add("[cyan]shapes[/cyan]")
    for s in master.shapes:
        shape_tree(s, shapes_node, show_xml=DUMP_XML)

    return root


# ── Slide ─────────────────────────────────────────────────────────────────────

def slide_tree(slide, idx: int) -> Tree:
    root = Tree(f"[bold green]slide[{idx}][/bold green]  layout={repr(slide.slide_layout.name)}")
    root.add(f"[cyan]shapes[/cyan] = {len(list(slide.shapes))}")
    shapes_node = root.add("[cyan]shapes[/cyan]")
    for s in slide.shapes:
        shape_tree(s, shapes_node, show_xml=DUMP_XML)

    # Notes
    try:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            root.add(f"[cyan]notes[/cyan] = {repr(notes_text[:200])}")
    except Exception:
        pass

    return root


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    if not args:
        console.print("[red]Usage: uv run python inspect_pptx.py <file.pptx> [--slides] [--xml][/red]")
        sys.exit(1)

    path = Path(args[0])
    if not path.exists():
        console.print(f"[red]File not found: {path}[/red]")
        sys.exit(1)

    prs = Presentation(str(path))

    # ── Presentation ──
    console.rule("[bold]Presentation[/bold]")
    t = Table("property", "value", box=box.SIMPLE_HEAVY, show_header=False)
    t.add_row("file", str(path))
    t.add_row("slide_width", emu(prs.slide_width))
    t.add_row("slide_height", emu(prs.slide_height))
    t.add_row("slides", str(len(prs.slides)))
    t.add_row("slide_layouts", str(len(prs.slide_layouts)))
    try:
        cp = prs.core_properties
        for field in ("title", "author", "subject", "keywords", "revision", "created", "modified"):
            t.add_row(f"core_properties.{field}", str(getattr(cp, field, None)))
    except Exception:
        pass
    console.print(t)

    # ── Slide master ──
    console.rule("[bold]Slide Master[/bold]")
    console.print(master_tree(prs.slide_master))

    # ── Layouts ──
    console.rule(f"[bold]Slide Layouts ({len(prs.slide_layouts)})[/bold]")
    for i, layout in enumerate(prs.slide_layouts):
        console.print(layout_tree(layout, i))
        console.print()

    # ── Slides (opt-in) ──
    if DUMP_SLIDES:
        console.rule(f"[bold]Slides ({len(prs.slides)})[/bold]")
        for i, slide in enumerate(prs.slides):
            console.print(slide_tree(slide, i))
            console.print()
    else:
        console.print(f"[dim](pass --slides to dump all {len(prs.slides)} content slides)[/dim]")


if __name__ == "__main__":
    main()
